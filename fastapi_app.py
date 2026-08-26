import hashlib
import json
import logging
import os
import posixpath
import re
import secrets
import smtplib
import uuid
from datetime import datetime, timedelta
from email.message import EmailMessage
from html.parser import HTMLParser
from io import BytesIO

from fastapi import Depends, FastAPI, File, Form, HTTPException, Request, UploadFile
from fastapi.responses import FileResponse, HTMLResponse, JSONResponse, RedirectResponse, Response
from fastapi.staticfiles import StaticFiles
from fastapi.templating import Jinja2Templates
from sqlalchemy import and_, func, or_
from sqlalchemy.orm import Session
from starlette.middleware.sessions import SessionMiddleware

import fastapi_config as cfg
from fastapi_auth import admin_from_request, hash_password, username_from_email, verify_password
from fastapi_db import (Admin, CertificateAward, Company, Course, Enrollment, Lesson, LessonMaterial, LessonProgress,
                        LearnerProfile, PasswordResetToken, Program, Purchase, Quiz, QuizAttempt, SessionObjective,
                        Settings, Student, db_session as next_db_session, ensure_schema, get_db)
from fastapi_storage import (guess_content_type, list_objects, object_bytes, object_key, package_object_key,
                             presigned_download_url, presigned_upload_url, r2_enabled, upload_fileobj)


CERTIFICATE_LEVEL_HOURS = 15
MODULES_PER_LEVEL = 5
SESSIONS_PER_MODULE = 5
SESSION_DURATION_MINUTES = 180
MASTER_CERTIFICATE_LEVEL = 4
PASSWORD_RESET_TOKEN_MINUTES = 60
MAX_LESSON_QUIZ_ATTEMPTS = 3
QUIZ_PASS_SCORE = 60
AI_CONTEXT_MAX_CHARS = 12000
AI_CONTEXT_MATERIAL_LIMIT = 8

MATERIAL_TYPE_OPTIONS = [
    ('slide', 'Slides / PPT / PDF'),
    ('html', 'HTML lessons / interactive slides'),
    ('module_intro_video', 'Module intro videos'),
    ('video', 'Videos'),
    ('case_application', 'Case applications'),
    ('case_study', 'Case studies'),
    ('case_analysis', 'Module case analysis'),
    ('toolkit', 'Toolkits'),
    ('toolkit_asset', 'Toolkit Excel/assets'),
    ('simulation', 'Simulations'),
    ('general_simulation', 'Module general simulation'),
    ('syllabus', 'Syllabus'),
    ('clo', 'CLO / CBO'),
    ('article', 'Articles / readings'),
    ('other', 'Other material'),
]
MATERIAL_TYPE_KEYS = {key for key, _label in MATERIAL_TYPE_OPTIONS}

EXPERTISE_AREAS = [
    {'name': 'AI Agents & Generative AI', 'slug': 'ai-agents-generative-ai'},
    {'name': 'Data Analytics & Business Intelligence', 'slug': 'data-analytics-bi'},
    {'name': 'Cybersecurity', 'slug': 'cybersecurity'},
    {'name': 'Cloud, DevOps & Infrastructure', 'slug': 'cloud-devops-infrastructure'},
    {'name': 'Automation & Workflow Design', 'slug': 'automation-workflow-design'},
    {'name': 'Digital Transformation Leadership', 'slug': 'digital-transformation-leadership'},
    {'name': 'Sustainability & Green Business', 'slug': 'sustainability-green-business'},
    {'name': 'Aviation Management & Operations', 'slug': 'aviation-management-operations'},
]

CERTIFICATE_LEVELS = [
    {'level': 1, 'name': 'Level 1 Certificate', 'hours': CERTIFICATE_LEVEL_HOURS, 'components': ['videos', 'practices', 'simulations', 'evaluation']},
    {'level': 2, 'name': 'Level 2 Certificate', 'hours': CERTIFICATE_LEVEL_HOURS, 'components': ['videos', 'practices', 'simulations', 'evaluation']},
    {'level': 3, 'name': 'Level 3 Certificate', 'hours': CERTIFICATE_LEVEL_HOURS, 'components': ['videos', 'practices', 'simulations', 'evaluation']},
    {'level': MASTER_CERTIFICATE_LEVEL, 'name': 'Master Certificate', 'hours': CERTIFICATE_LEVEL_HOURS * 3, 'components': ['Level 1 completion', 'Level 2 completion', 'Level 3 completion']},
]


cfg.validate_production()
app = FastAPI(title='Hub Academy')
app.add_middleware(
    SessionMiddleware,
    secret_key=cfg.SECRET_KEY,
    same_site='lax',
    https_only=os.environ.get('ENV') == 'production' or os.environ.get('APP_ENV') == 'production',
)
app.mount('/static', StaticFiles(directory=os.path.join(cfg.BASE_DIR, 'static')), name='static')

templates = Jinja2Templates(directory=os.path.join(cfg.BASE_DIR, 'fastapi_templates'))
logger = logging.getLogger(__name__)

# Additive: certificate PDF download + public verification (self-contained module).
import certificate_verify  # noqa: E402
app.include_router(certificate_verify.router)


@app.on_event('startup')
def startup():
    ensure_schema()


@app.get('/healthz')
def healthz():
    return {'ok': True}


@app.get('/_build-check')
def build_check():
    lesson_template = os.path.join(cfg.BASE_DIR, 'fastapi_templates', 'learn', 'lesson.html')
    try:
        with open(lesson_template, 'r', encoding='utf-8') as handle:
            text = handle.read()
    except OSError:
        text = ''
    return {
        'commit_marker': 'aef1119-fa5ca6e',
        'lesson_workspace': 'lesson-workspace' in text,
        'ai_tools': 'AI Tools & Applications' in text,
        'self_evaluation': 'Self-Evaluation' in text,
        'service_worker_v2': True,
    }


@app.get('/manifest.webmanifest')
def pwa_manifest():
    return FileResponse(os.path.join(cfg.BASE_DIR, 'static', 'manifest.webmanifest'), media_type='application/manifest+json')


@app.get('/service-worker.js')
def service_worker():
    return FileResponse(os.path.join(cfg.BASE_DIR, 'static', 'service-worker.js'), media_type='application/javascript')


def slugify(value):
    value = re.sub(r'[^a-zA-Z0-9]+', '-', value.lower()).strip('-')
    return value or 'course'


def course_slug(course):
    return course.slug or f'{course.id}-{slugify(course.title)}'


def course_price(course):
    amount = course.price_cents or 0
    if amount <= 0 or course.allow_free_enrollment:
        return 'Free'
    return f"{(course.currency or 'USD').upper()} {amount / 100:.2f}"


COURSE_IMAGE_RULES = [
    (('cyber', 'security', 'privacy', 'risk'), [
        'https://images.unsplash.com/photo-1563986768609-322da13575f3?auto=format&fit=crop&w=900&q=80',
        'https://images.unsplash.com/photo-1516321318423-f06f85e504b3?auto=format&fit=crop&w=900&q=80',
        'https://images.unsplash.com/photo-1550751827-4bd374c3f58b?auto=format&fit=crop&w=900&q=80',
        'https://images.unsplash.com/photo-1510511459019-5dda7724fd87?auto=format&fit=crop&w=900&q=80',
    ]),
    (('ai', 'artificial intelligence', 'machine learning', 'readiness', 'automation'), [
        'https://images.unsplash.com/photo-1677442136019-21780ecad995?auto=format&fit=crop&w=900&q=80',
        'https://images.unsplash.com/photo-1677756119517-756a188d2d94?auto=format&fit=crop&w=900&q=80',
        'https://images.unsplash.com/photo-1620712943543-bcc4688e7485?auto=format&fit=crop&w=900&q=80',
        'https://images.unsplash.com/photo-1639322537228-f710d846310a?auto=format&fit=crop&w=900&q=80',
        'https://images.unsplash.com/photo-1555255707-c07966088b7b?auto=format&fit=crop&w=900&q=80',
    ]),
    (('data', 'analytics', 'dashboard', 'intelligence', 'metrics'), [
        'https://images.unsplash.com/photo-1551288049-bebda4e38f71?auto=format&fit=crop&w=900&q=80',
        'https://images.unsplash.com/photo-1460925895917-afdab827c52f?auto=format&fit=crop&w=900&q=80',
        'https://images.unsplash.com/photo-1553877522-43269d4ea984?auto=format&fit=crop&w=900&q=80',
        'https://images.unsplash.com/photo-1504868584819-f8e8b4b6d7e3?auto=format&fit=crop&w=900&q=80',
    ]),
    (('digital', 'maturity', 'transformation', 'platform'), [
        'https://images.unsplash.com/photo-1551434678-e076c223a692?auto=format&fit=crop&w=900&q=80',
        'https://images.unsplash.com/photo-1517245386807-bb43f82c33c4?auto=format&fit=crop&w=900&q=80',
        'https://images.unsplash.com/photo-1556761175-b413da4baf72?auto=format&fit=crop&w=900&q=80',
    ]),
    (('cloud', 'infrastructure', 'architecture'), [
        'https://images.unsplash.com/photo-1451187580459-43490279c0fa?auto=format&fit=crop&w=900&q=80',
        'https://images.unsplash.com/photo-1497366754035-f200968a6e72?auto=format&fit=crop&w=900&q=80',
        'https://images.unsplash.com/photo-1518770660439-4636190af475?auto=format&fit=crop&w=900&q=80',
    ]),
    (('finance', 'financial', 'cost', 'budget'), [
        'https://images.unsplash.com/photo-1460925895917-afdab827c52f?auto=format&fit=crop&w=900&q=80',
        'https://images.unsplash.com/photo-1554224155-6726b3ff858f?auto=format&fit=crop&w=900&q=80',
        'https://images.unsplash.com/photo-1554224154-26032ffc0d07?auto=format&fit=crop&w=900&q=80',
    ]),
    (('leadership', 'strategy', 'change', 'governance'), [
        'https://images.unsplash.com/photo-1521737604893-d14cc237f11d?auto=format&fit=crop&w=900&q=80',
        'https://images.unsplash.com/photo-1556761175-4b46a572b786?auto=format&fit=crop&w=900&q=80',
    ]),
    (('sustainability', 'green', 'esg'), [
        'https://images.unsplash.com/photo-1497435334941-8c899ee9e8e9?auto=format&fit=crop&w=900&q=80',
        'https://images.unsplash.com/photo-1473341304170-971dccb5ac1e?auto=format&fit=crop&w=900&q=80',
        'https://images.unsplash.com/photo-1466611653911-95081537e5b7?auto=format&fit=crop&w=900&q=80',
    ]),
    (('project', 'process', 'operations'), [
        'https://images.unsplash.com/photo-1517245386807-bb43f82c33c4?auto=format&fit=crop&w=900&q=80',
        'https://images.unsplash.com/photo-1454165804606-c3d57bc86b40?auto=format&fit=crop&w=900&q=80',
        'https://images.unsplash.com/photo-1504384308090-c894fdcc538d?auto=format&fit=crop&w=900&q=80',
    ]),
]

DEFAULT_COURSE_IMAGES = [
    'https://images.unsplash.com/photo-1556761175-b413da4baf72?auto=format&fit=crop&w=900&q=80',
    'https://images.unsplash.com/photo-1522202176988-66273c2fd55f?auto=format&fit=crop&w=900&q=80',
    'https://images.unsplash.com/photo-1557804506-669a67965ba0?auto=format&fit=crop&w=900&q=80',
    'https://images.unsplash.com/photo-1542744173-8e7e53415bb0?auto=format&fit=crop&w=900&q=80',
]

# Title-specific bundled images (always load; meaningful; no recycling).
COURSE_IMAGE_LOCAL = [
    (('agentic', 'ai agent'), '/static/images/course-ai-agents.jpg'),
    (('engineering leadership',), '/static/images/course-engineering.jpg'),
    (('ai leadership', 'ai strategy'), '/static/images/course-ai-leadership.jpg'),
]


def deterministic_choice(items, seed):
    digest = hashlib.sha256(seed.encode('utf-8')).hexdigest()
    return items[int(digest[:8], 16) % len(items)]


def course_image_url(image_url, seed):
    signature = hashlib.sha256(seed.encode('utf-8')).hexdigest()[:10]
    separator = '&' if '?' in image_url else '?'
    return f'{image_url}{separator}sig={signature}'


def course_image(course):
    if course.thumbnail_url:
        url = course.thumbnail_url
        # Pasted URLs and local static paths are used as-is; an R2 object key is
        # served (presigned) through our own /courses/<id>/thumbnail route.
        if url.startswith(('http://', 'https://', '/')):
            return url
        return f'/courses/{course.id}/thumbnail'
    title_hay = ((course.title or '') + ' ' + (course.expertise_area or '')).lower()
    for keywords, local_image in COURSE_IMAGE_LOCAL:
        if any(keyword in title_hay for keyword in keywords):
            return local_image
    haystack = ' '.join([
        course.title or '',
        course.expertise_area or '',
        course.level or '',
        course.description or '',
        course.sales_copy or '',
    ]).lower()
    seed = f'{course.id or ""}|{course.title or ""}|{course.expertise_area or ""}|{course.certificate_level or ""}'
    for keywords, image_urls in COURSE_IMAGE_RULES:
        if any(keyword in haystack for keyword in keywords):
            return course_image_url(deterministic_choice(image_urls, seed), seed)
    return course_image_url(deterministic_choice(DEFAULT_COURSE_IMAGES, seed), seed)


def module_number_for_session(lesson_number: int):
    return int((max(lesson_number, 1) - 1) / SESSIONS_PER_MODULE) + 1


def session_number_for_lesson(lesson_number: int):
    return ((max(lesson_number, 1) - 1) % SESSIONS_PER_MODULE) + 1


def lesson_module_number(lesson):
    return lesson.module_number or module_number_for_session(lesson.lesson_number)


def lesson_session_number(lesson):
    return lesson.session_number or session_number_for_lesson(lesson.lesson_number)


def lesson_duration_minutes(lesson):
    return lesson.duration_minutes or SESSION_DURATION_MINUTES


def lesson_has_content(db: Session, lesson_id: int) -> bool:
    return bool(
        db.query(LessonMaterial).filter_by(lesson_id=lesson_id).count()
        or db.query(LessonProgress).filter_by(lesson_id=lesson_id).count()
        or db.query(Quiz).filter_by(lesson_id=lesson_id).count()
    )


def normalize_course_modules(db: Session, course: Course, module_count: int = MODULES_PER_LEVEL):
    lessons = db.query(Lesson).filter_by(course_id=course.id).order_by(Lesson.lesson_number).all()
    for number in range(len(lessons) + 1, module_count + 1):
        lesson = Lesson(course_id=course.id, lesson_number=number, created_at=datetime.utcnow())
        db.add(lesson)
        lessons.append(lesson)
    blocked_extras = []
    for index, lesson in enumerate(lessons, start=1):
        if index <= module_count:
            lesson.lesson_number = index
            lesson.module_number = index
            lesson.session_number = 1
            lesson.duration_minutes = SESSION_DURATION_MINUTES
            if re.match(r'^Module\s+\d+(\s*[-–]\s*Session\s+\d+)?$', lesson.title or '', flags=re.IGNORECASE):
                lesson.title = f'Module {index}'
            elif not (lesson.title or '').strip():
                lesson.title = f'Module {index}'
        elif lesson_has_content(db, lesson.id):
            blocked_extras.append(lesson.title or f'Lesson {lesson.lesson_number}')
        else:
            db.delete(lesson)
    if blocked_extras:
        return blocked_extras
    course.num_lessons = module_count
    return []


def material_display_label(material, index):
    material_type = material.material_type or 'other'
    if material_type in ('html', 'slide'):
        return f'Session {index}'
    if material_type == 'module_intro_video':
        return 'Module intro video'
    if material_type == 'video':
        return f'Video {index}'
    if material_type == 'simulation':
        return f'Simulation {index}'
    if material_type == 'general_simulation':
        return 'General simulation'
    if material_type == 'case_analysis':
        return 'Case analysis'
    if material_type in ('toolkit', 'toolkit_asset'):
        return f'Application {index}'
    if material_type in ('case_application', 'case_study'):
        return f'Case {index}'
    if material_type == 'syllabus':
        return 'Syllabus'
    if material_type == 'clo':
        return 'CLO / CBO'
    return f'Material {index}'


def group_lesson_materials(materials):
    groups = {
        'learning': [],
        'ai_tools': [],
        'support': [],
    }
    counts = {}
    for material in materials:
        material_type = material.material_type or 'other'
        counts[material_type] = counts.get(material_type, 0) + 1
        item = {
            'material': material,
            'label': material_display_label(material, counts[material_type]),
        }
        if material_type in ('html', 'slide', 'video', 'simulation'):
            groups['learning'].append(item)
        elif material_type in ('toolkit', 'toolkit_asset', 'case_application', 'case_study'):
            groups['ai_tools'].append(item)
        else:
            groups['support'].append(item)
    return groups


def module_material_groups(materials, module_number):
    groups = {
        'sessions': [],
        'videos': [],
        'simulations': [],
        'applications': [],
        'cases': [],
        'references': [],
    }
    counters = {key: 0 for key in groups}
    for material in materials:
        material_type = material.material_type or 'other'
        if material_type == 'toolkit_asset':
            continue
        if material_type in ('html', 'slide'):
            group = 'sessions'
            label = f'Session {module_number}'
        elif material_type == 'video':
            group = 'videos'
            label = f'Video {module_number}'
        elif material_type == 'simulation':
            group = 'simulations'
            label = f'Simulation {module_number}'
        elif material_type == 'toolkit' and html_material(material):
            group = 'applications'
            label = f'Application {module_number}'
        elif material_type in ('case_application', 'case_study'):
            group = 'cases'
            label = f'Case {module_number}'
        else:
            group = 'references'
            counters[group] += 1
            label = material_display_label(material, counters[group])
        groups[group].append({'material': material, 'label': label})
    return groups


def material_unit_number(material):
    file_text = ' '.join([
        material.file_name or '',
        material.video_name or '',
    ])
    path_text = ' '.join([
        material.object_key or '',
        material.file_path or '',
    ])
    text = f'{file_text} {path_text}'

    folder_match = re.search(r'(?:^|[/\\])0*(\d+)[-_ ]session(?:[_-]\d+min)?(?:[/\\]|$)', text, flags=re.IGNORECASE)
    if folder_match:
        try:
            number = int(folder_match.group(1))
            if number > 0:
                return number
        except ValueError:
            pass

    filename_patterns = [
        r'(?:session|simulation|sim|vol|application|case)[-_ ]*0*(\d+)',
        r'(?:^|[_\W])module[-_ ]*0*(\d+)',
        r'[_-]0*(\d+)[_-](?:ar[_-])?audio',
    ]
    for pattern in filename_patterns:
        match = re.search(pattern, file_text, flags=re.IGNORECASE)
        if match:
            try:
                number = int(match.group(1))
                if number > 0:
                    return number
            except ValueError:
                continue

    # Bulk imports store packages such as module-42-session-3-html.
    # This must not treat 01_Session_30min as session 30.
    session_match = re.search(r'(?:^|[/_-])session[-_ ]*0*(\d+)(?:[-_/]|$)', path_text, flags=re.IGNORECASE)
    if session_match:
        try:
            number = int(session_match.group(1))
            if number > 0:
                return number
        except ValueError:
            pass

    fallback_patterns = [
        r'(?:session|simulation|sim|vol|application|case)[-_ ]*0*(\d+)',
        r'(?:^|[_\W])module[-_ ]*0*(\d+)',
        r'[_-]0*(\d+)[_-](?:ar[_-])?audio',
        r'[_-]0*(\d+)[_-][a-z]',
        r'[_-]0*(\d+)\.',
    ]
    for pattern in fallback_patterns:
        match = re.search(pattern, text, flags=re.IGNORECASE)
        if match:
            try:
                number = int(match.group(1))
                if number > 0:
                    return number
            except ValueError:
                continue
    return None


def module_session_groups(materials, module_number, objective_map=None):
    objective_map = objective_map or {}
    sessions = {}
    fallback_counts = {'session': 0, 'video': 0, 'simulation': 0, 'application': 0, 'case': 0}

    def session_entry(unit):
        sessions.setdefault(unit, {
            'unit': unit,
            'label': f'{module_number}.{unit}',
            'objective': objective_map.get((module_number, unit)) or 'Review the session resources, apply the activity, and prepare for the module evaluation.',
            'session': [],
            'video': [],
            'simulation': [],
            'application': [],
            'case': [],
        })
        return sessions[unit]

    for material in materials:
        material_type = material.material_type or 'other'
        if material_type in ('toolkit_asset', 'module_intro_video', 'case_analysis', 'general_simulation'):
            continue
        if material_type in ('html', 'slide'):
            bucket = 'session'
        elif material_type == 'video':
            bucket = 'video'
        elif material_type == 'simulation':
            bucket = 'simulation'
        elif material_type == 'toolkit' and html_material(material):
            bucket = 'application'
        elif material_type in ('case_application', 'case_study'):
            bucket = 'case'
        else:
            continue
        unit = material_unit_number(material)
        if not unit:
            fallback_counts[bucket] += 1
            unit = fallback_counts[bucket]
        session_entry(unit)[bucket].append(material)

    return [sessions[key] for key in sorted(sessions)]


def module_extra_materials(materials):
    extras = {
        'module_intro_video': [],
        'case_analysis': [],
        'general_simulation': [],
        'references': [],
    }
    for material in materials:
        material_type = material.material_type or 'other'
        if material_type in extras:
            extras[material_type].append(material)
        elif material_type in ('syllabus', 'clo'):
            extras['references'].append(material)
    return extras


def session_objective_map(db: Session, course_id: int):
    rows = db.query(SessionObjective).filter_by(course_id=course_id).all()
    return {
        (row.module_number, row.session_number): row.objective
        for row in rows
        if row.objective
    }


def current_module_number_for_lesson(db: Session, lesson: Lesson) -> int:
    lessons = db.query(Lesson).filter_by(course_id=lesson.course_id).order_by(Lesson.lesson_number).limit(lesson.course.num_lessons or MODULES_PER_LEVEL).all()
    return next((index for index, item in enumerate(lessons, start=1) if item.id == lesson.id), lesson_session_number(lesson))


def extract_pdf_text(pdf_bytes: bytes) -> str:
    try:
        from pypdf import PdfReader
    except Exception as exc:
        raise RuntimeError('pypdf is required for syllabus objective extraction.') from exc
    reader = PdfReader(BytesIO(pdf_bytes))
    return '\n'.join(page.extract_text() or '' for page in reader.pages)


def extract_topic_objectives(text: str, limit: int = SESSIONS_PER_MODULE):
    lines = [line.strip() for line in (text or '').splitlines()]
    start = next((index for index, line in enumerate(lines) if re.search(r'\btopics covered\b', line, re.IGNORECASE)), -1)
    if start < 0:
        return []
    objectives = []
    stop_pattern = re.compile(r'^(applied lab|module\s+\d+|learning outcomes|note on scope)\b', re.IGNORECASE)
    for line in lines[start + 1:]:
        cleaned = re.sub(r'^[\u2022\-–—*\d\.\)\s]+', '', line).strip()
        if not cleaned:
            continue
        if stop_pattern.search(cleaned):
            break
        if len(cleaned) < 8:
            continue
        objectives.append(cleaned)
        if len(objectives) >= limit:
            break
    return objectives


def save_session_objectives(db: Session, course_id: int, module_number: int, objectives):
    saved = 0
    for index, objective in enumerate(objectives[:SESSIONS_PER_MODULE], start=1):
        existing = db.query(SessionObjective).filter_by(
            course_id=course_id,
            module_number=module_number,
            session_number=index,
        ).first()
        if not existing:
            existing = SessionObjective(course_id=course_id, module_number=module_number, session_number=index)
            db.add(existing)
        existing.objective = objective.strip()
        existing.title = None
        existing.source = 'syllabus'
        saved += 1
    db.commit()
    return saved


def module_nav_for_lessons(db: Session, student_id: int, lessons, current_lesson_id: int):
    items = []
    previous_progress = None
    for index, item in enumerate(lessons, start=1):
        progress = db.query(LessonProgress).filter_by(student_id=student_id, lesson_id=item.id).first()
        unlocked = index == 1 or session_unlocked(previous_progress)
        items.append({
            'lesson': item,
            'number': index,
            'active': item.id == current_lesson_id,
            'unlocked': unlocked,
            'completed': bool(progress and progress.is_completed),
        })
        previous_progress = progress
    return items


def lesson_context_text(lesson, materials):
    course = lesson.course
    material_lines = []
    for material in materials[:40]:
        label = material.material_type.replace('_', ' ').title()
        if material.video_url:
            material_lines.append(f"- {label}: external video link")
        elif material.content_type:
            material_lines.append(f"- {label}: {material.content_type}")
        else:
            material_lines.append(f"- {label}")
    return "\n".join([
        f"Course: {course.title if course else ''}",
        f"Lesson: {lesson.title}",
        f"Description: {lesson.description or ''}",
        "Available material types:",
        "\n".join(material_lines),
    ])


class ReadableHTMLParser(HTMLParser):
    def __init__(self):
        super().__init__()
        self.parts = []
        self.skip_depth = 0

    def handle_starttag(self, tag, attrs):
        if tag in {'script', 'style', 'noscript', 'svg'}:
            self.skip_depth += 1
        if tag in {'p', 'div', 'section', 'article', 'li', 'br', 'h1', 'h2', 'h3', 'h4', 'tr'}:
            self.parts.append('\n')

    def handle_endtag(self, tag):
        if tag in {'script', 'style', 'noscript', 'svg'} and self.skip_depth:
            self.skip_depth -= 1
        if tag in {'p', 'div', 'section', 'article', 'li', 'h1', 'h2', 'h3', 'h4', 'tr'}:
            self.parts.append('\n')

    def handle_data(self, data):
        if not self.skip_depth:
            value = re.sub(r'\s+', ' ', data or '').strip()
            if value:
                self.parts.append(value)

    def text(self):
        return re.sub(r'\n{3,}', '\n\n', '\n'.join(self.parts)).strip()


def decode_text_bytes(raw_bytes):
    for encoding in ('utf-8', 'utf-8-sig', 'cp1252', 'latin-1'):
        try:
            return raw_bytes.decode(encoding)
        except UnicodeDecodeError:
            continue
    return raw_bytes.decode('utf-8', errors='ignore')


def html_to_text(raw_html):
    parser = ReadableHTMLParser()
    parser.feed(raw_html)
    return parser.text()


def extract_pptx_text(raw_bytes):
    from pptx import Presentation
    presentation = Presentation(BytesIO(raw_bytes))
    parts = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        slide_parts = []
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text:
                slide_parts.append(shape.text)
        if slide_parts:
            parts.append(f"Slide {slide_index}: " + ' '.join(slide_parts))
    return '\n'.join(parts)


def extract_xlsx_text(raw_bytes):
    from openpyxl import load_workbook
    workbook = load_workbook(BytesIO(raw_bytes), read_only=True, data_only=True)
    parts = []
    for sheet in workbook.worksheets[:4]:
        rows = []
        for row in sheet.iter_rows(max_row=40, values_only=True):
            values = [str(value).strip() for value in row if value is not None and str(value).strip()]
            if values:
                rows.append(' | '.join(values))
        if rows:
            parts.append(f"Sheet {sheet.title}: " + ' '.join(rows))
    return '\n'.join(parts)


def material_source_bytes(material, max_bytes=20_000_000):
    if material.size_bytes and material.size_bytes > max_bytes:
        return None
    key = material.object_key or material.file_path
    if key:
        try:
            return object_bytes(key)
        except Exception as exc:
            logger.warning('Could not read material %s from R2: %s', material.id, exc)
            return None
    if material.file_path and os.path.exists(material.file_path):
        try:
            with open(material.file_path, 'rb') as handle:
                return handle.read(max_bytes)
        except Exception as exc:
            logger.warning('Could not read local material %s: %s', material.id, exc)
    return None


def extract_material_text(material):
    material_type = (material.material_type or '').lower()
    content_type = (material.content_type or '').lower()
    name = (material.file_name or material.object_key or material.file_path or '').lower()
    if material_type in {'video', 'module_intro_video'} or content_type.startswith('video/'):
        return ''
    if not any([
        'html' in content_type,
        'text/' in content_type,
        'json' in content_type,
        'pdf' in content_type,
        'presentation' in content_type,
        'spreadsheet' in content_type,
        name.endswith(('.html', '.htm', '.txt', '.md', '.json', '.pdf', '.pptx', '.xlsx')),
    ]):
        return ''
    raw_bytes = material_source_bytes(material)
    if not raw_bytes:
        return ''
    try:
        if 'pdf' in content_type or name.endswith('.pdf'):
            return extract_pdf_text(raw_bytes)
        if 'presentation' in content_type or name.endswith('.pptx'):
            return extract_pptx_text(raw_bytes)
        if 'spreadsheet' in content_type or name.endswith('.xlsx'):
            return extract_xlsx_text(raw_bytes)
        raw_text = decode_text_bytes(raw_bytes)
        if 'html' in content_type or name.endswith(('.html', '.htm')):
            return html_to_text(raw_text)
        return re.sub(r'\s+', ' ', raw_text).strip()
    except Exception as exc:
        logger.warning('Could not extract material %s text: %s', material.id, exc)
        return ''


def lesson_learning_corpus(lesson, materials):
    sections = []
    total_chars = 0
    readable_materials = [
        material for material in materials
        if (material.material_type or '').lower() in {
            'html', 'slide', 'simulation', 'toolkit', 'case_application', 'case_analysis',
            'general_simulation', 'syllabus', 'clo', 'article', 'other'
        }
    ]
    for material in readable_materials[:AI_CONTEXT_MATERIAL_LIMIT]:
        text = extract_material_text(material)
        if not text:
            continue
        text = re.sub(r'\s+', ' ', text).strip()
        if len(text) < 120:
            continue
        remaining = AI_CONTEXT_MAX_CHARS - total_chars
        if remaining <= 0:
            break
        excerpt = text[:remaining]
        total_chars += len(excerpt)
        label = (material.material_type or 'material').replace('_', ' ').title()
        sections.append(f"{label} content:\n{excerpt}")
    if sections:
        return "\n\n".join(sections)
    return ''


def openai_api_key(db: Session = None, course=None):
    if course and getattr(course, 'openai_api_key_override', None):
        return course.openai_api_key_override.strip()
    if cfg.OPENAI_API_KEY:
        return cfg.OPENAI_API_KEY.strip()
    if db:
        setting = db.query(Settings).filter_by(key='openai_api_key').first()
        if setting and setting.value:
            return setting.value.strip()
    return ''


def openai_client(db: Session = None, course=None):
    api_key = openai_api_key(db, course)
    if not api_key:
        raise RuntimeError('OpenAI API key is not configured.')
    from openai import OpenAI
    return OpenAI(api_key=api_key, timeout=45, max_retries=1)


def parse_json_response(text, fallback):
    cleaned = (text or '').strip()
    if cleaned.startswith('```json'):
        cleaned = cleaned[7:]
    if cleaned.startswith('```'):
        cleaned = cleaned[3:]
    if cleaned.endswith('```'):
        cleaned = cleaned[:-3]
    try:
        return json.loads(cleaned.strip())
    except Exception:
        return fallback


def ai_error_response(exc, fallback_message):
    message = str(exc)
    if 'OpenAI API key is not configured' in message:
        return JSONResponse({'error': 'OpenAI API key is not configured. Add it in Admin Settings or Cloud Run OPENAI_API_KEY.'}, status_code=400)
    logger.exception('%s: %s', fallback_message, exc)
    details = []
    current = exc
    while current and len(details) < 4:
        detail = str(current).strip() or current.__class__.__name__
        details.append(detail.replace('\n', ' '))
        current = getattr(current, '__cause__', None) or getattr(current, '__context__', None)
    public_detail = ' | '.join(dict.fromkeys(details))
    if len(public_detail) > 260:
        public_detail = public_detail[:257] + '...'
    return JSONResponse({'error': f'{fallback_message} OpenAI said: {public_detail}' if public_detail else fallback_message}, status_code=502)


def openai_chat_completion(client, messages, max_tokens, temperature):
    models = []
    for model in [cfg.OPENAI_MODEL, 'gpt-4o-mini', 'gpt-4.1-mini']:
        if model and model not in models:
            models.append(model)
    last_error = None
    for model in models:
        try:
            return client.chat.completions.create(
                model=model,
                messages=messages,
                max_tokens=max_tokens,
                temperature=temperature,
            )
        except Exception as exc:
            last_error = exc
    raise last_error or RuntimeError('OpenAI chat completion failed.')


def ai_generate_quiz_questions(db, lesson, materials, previous_questions):
    corpus = lesson_learning_corpus(lesson, materials)
    prompt = f"""Create exactly 5 multiple-choice questions for this lesson.

The questions must be scenario-based, non-obvious, and specific to the learning content below.
Never ask about file names, document names, upload structure, folders, or course platform mechanics.
If the content is thin, ask about concepts, decisions, risks, examples, and application from the content.

Lesson context:
{lesson_context_text(lesson, materials)}

Learning content extracted from materials:
{corpus or 'No readable material text was extracted. Use the lesson description only, and do not mention files.'}

Previous questions to avoid repeating:
{json.dumps(previous_questions[-30:], ensure_ascii=False)}

Return only valid JSON:
[
  {{
    "question": "question text",
    "options": ["A", "B", "C", "D"],
    "correct_answer": 0,
    "explanation": "short explanation"
  }}
]"""
    response = openai_chat_completion(
        openai_client(db, lesson.course),
        [{'role': 'user', 'content': prompt}],
        max_tokens=1800,
        temperature=0.9,
    )
    data = parse_json_response(response.choices[0].message.content, [])
    return data[:5] if isinstance(data, list) else []


def ai_generate_feedback(db, lesson, score, questions, answers):
    prompt = f"""A learner just completed a quiz for this lesson.

Lesson: {lesson.title}
Score: {score}%
Questions and correct answers:
{json.dumps(questions, ensure_ascii=False)}
Learner answers by question index:
{json.dumps(answers, ensure_ascii=False)}

Address the learner directly as "You". Be concise and encouraging.
Return only valid JSON with two short bullet lists. Each bullet is ONE short sentence (max ~14 words), starting with "You".
Give 2-3 strengths and 2-3 improvements (fewer if the quiz was short).
{{
  "strengths": ["You ...", "You ..."],
  "improvements": ["You ...", "You ..."]
}}"""
    response = openai_chat_completion(
        openai_client(db, lesson.course),
        [{'role': 'user', 'content': prompt}],
        max_tokens=500,
        temperature=0.5,
    )
    data = parse_json_response(response.choices[0].message.content, {})

    def clean_bullets(value):
        if isinstance(value, str):
            value = [value]
        if not isinstance(value, list):
            return []
        return [str(item).strip() for item in value if str(item).strip()][:3]

    strengths = clean_bullets(data.get('strengths')) or ['You completed the quiz — nice work.']
    improvements = clean_bullets(data.get('improvements')) or ['You should review the lesson materials before your next attempt.']
    return {
        'strengths': strengths,
        'improvements': improvements,
        # joined text kept for the stored attempt history
        'feedback': ' • '.join(strengths),
        'recommendations': ' • '.join(improvements),
    }


def course_quiz_stats(db: Session, student_id: int, course_id: int):
    attempts = db.query(QuizAttempt).join(Quiz, QuizAttempt.quiz_id == Quiz.id).join(
        Lesson, Quiz.lesson_id == Lesson.id
    ).filter(QuizAttempt.student_id == student_id, Lesson.course_id == course_id).all()
    if not attempts:
        return {'attempts': 0, 'average': 0, 'best': 0}
    scores = [attempt.score for attempt in attempts]
    return {'attempts': len(scores), 'average': int(sum(scores) / len(scores)), 'best': max(scores)}


def course_quiz_average_passed(db: Session, student_id: int, course_id: int):
    stats = course_quiz_stats(db, student_id, course_id)
    return stats['attempts'] > 0 and stats['average'] >= QUIZ_PASS_SCORE


def certificate_badge(course):
    level = course.certificate_level or 0
    if level in (1, 2, 3):
        hours = course.learning_hours or CERTIFICATE_LEVEL_HOURS
        return f'Level {level} certificate track - {hours} hours'
    return 'Certificate-ready'


def step_state(progress):
    return [
        {
            'number': 1,
            'key': 'learn',
            'title': 'Learn',
            'label': 'Watch and study',
            'complete': bool(progress and progress.content_viewed),
            'unlocked': True,
        },
        {
            'number': 2,
            'key': 'practice',
            'title': 'Practice',
            'label': 'Apply the concept',
            'complete': bool(progress and progress.revise_viewed),
            'unlocked': bool(progress and progress.content_viewed),
        },
        {
            'number': 3,
            'key': 'validate',
            'title': 'Validate',
            'label': 'Check and confirm',
            'complete': bool(progress and progress.quiz_completed),
            'unlocked': bool(progress and progress.content_viewed and progress.revise_viewed),
        },
    ]


def session_unlocked(previous_progress):
    # A module unlocks only when the previous module is completed.
    # (A missing progress row means the previous module was never finished -> stay locked.)
    return bool(previous_progress and previous_progress.is_completed)


def journey_for_course(db: Session, course: Course, student_id: int):
    lessons = db.query(Lesson).filter_by(course_id=course.id).order_by(Lesson.lesson_number).limit(course.num_lessons or MODULES_PER_LEVEL).all()
    progress = {p.lesson_id: p for p in db.query(LessonProgress).filter(
        LessonProgress.student_id == student_id,
        LessonProgress.lesson_id.in_([lesson.id for lesson in lessons] or [0])
    ).all()}
    modules = []
    previous_progress = None
    total_done = 0
    for idx, lesson in enumerate(lessons):
        current_progress = progress.get(lesson.id)
        if current_progress and current_progress.is_completed:
            total_done += 1
        module_number = lesson_module_number(lesson)
        while len(modules) < module_number:
            modules.append({'number': len(modules) + 1, 'sessions': []})
        modules[module_number - 1]['sessions'].append({
            'lesson': lesson,
            'progress': current_progress,
            'unlocked': idx == 0 or session_unlocked(previous_progress),
            'session_number': lesson_session_number(lesson),
            'duration_minutes': lesson_duration_minutes(lesson),
            'steps': step_state(current_progress),
        })
        previous_progress = current_progress
    return {
        'modules': modules,
        'total_sessions': len(lessons),
        'completed_sessions': total_done,
        'progress': progress,
        'pct': int(total_done / len(lessons) * 100) if lessons else 0,
    }


def subscription_plans():
    return [
        {'name': 'Monthly', 'price': cfg.SUBSCRIPTION_MONTHLY_PRICE, 'period': 'month', 'stripe_price_id': cfg.STRIPE_MONTHLY_PRICE_ID,
         'features': ['Unlimited course access', 'AI learning guide', 'Progress tracking', 'Cancel anytime']},
        {'name': 'Annual', 'price': cfg.SUBSCRIPTION_ANNUAL_PRICE, 'period': 'year', 'stripe_price_id': cfg.STRIPE_ANNUAL_PRICE_ID,
         'features': ['Everything in Monthly', 'Best value', 'Certificates included', 'Priority support']},
    ]


def student_from_request(request: Request, db: Session):
    student_id = request.session.get('student_id')
    if not student_id:
        return None
    student = db.get(Student, int(student_id))
    if not student or not student.is_active:
        request.session.clear()
        return None
    return student


def require_admin(request: Request, db: Session):
    admin = admin_from_request(request, db)
    if not admin:
        raise HTTPException(status_code=303, headers={'Location': '/admin/login'})
    return admin


def template(request: Request, name: str, db: Session, context=None):
    ctx = {
        'request': request,
        'current_user': student_from_request(request, db),
        'current_admin': admin_from_request(request, db),
        'course_slug': course_slug,
        'course_price': course_price,
        'course_image': course_image,
        'certificate_badge': certificate_badge,
        'lesson_module_number': lesson_module_number,
        'lesson_session_number': lesson_session_number,
        'lesson_duration_minutes': lesson_duration_minutes,
        'certificate_levels': CERTIFICATE_LEVELS,
        'expertise_areas': EXPERTISE_AREAS,
        'plans': subscription_plans(),
    }
    if context:
        ctx.update(context)
    return templates.TemplateResponse(name, ctx)


def categories(db: Session):
    rows = db.query(Course.expertise_area, func.count(Course.id)).filter(Course.is_published.is_(True)).group_by(Course.expertise_area).all()
    counts = {name: count for name, count in rows if name}
    return [
        {'name': area['name'], 'count': counts.get(area['name']), 'href': f"/courses?expertise={area['name']}"}
        for area in EXPERTISE_AREAS
    ]


def find_course(db: Session, identifier: str):
    course = db.query(Course).filter_by(slug=identifier, is_published=True).first()
    if course:
        return course
    prefix = identifier.split('-', 1)[0]
    if prefix.isdigit():
        return db.query(Course).filter_by(id=int(prefix), is_published=True).first()
    return None


def enroll_student(db: Session, student_id: int, course_id: int):
    enrollment = db.query(Enrollment).filter_by(student_id=student_id, course_id=course_id).first()
    if enrollment:
        enrollment.is_active = True
    else:
        db.add(Enrollment(student_id=student_id, course_id=course_id, is_active=True))


def material_access_allowed(material: LessonMaterial, request: Request, db: Session) -> bool:
    admin = admin_from_request(request, db)
    if admin:
        return True
    student = student_from_request(request, db)
    if student and material.lesson:
        return db.query(Enrollment).filter_by(
            student_id=student.id,
            course_id=material.lesson.course_id,
            is_active=True,
        ).first() is not None
    return False


def html_material(material: LessonMaterial) -> bool:
    name = (material.file_name or '').lower()
    content_type = (material.content_type or '').lower()
    return name.endswith(('.html', '.htm')) or content_type.startswith('text/html')


def uploaded_file_size(fileobj) -> int:
    try:
        current = fileobj.tell()
        fileobj.seek(0, os.SEEK_END)
        size = fileobj.tell()
        fileobj.seek(current)
        return int(size)
    except Exception:
        return 0


def infer_material_type_from_key(key: str) -> str:
    parts = key.split('/')
    if len(parts) >= 3 and parts[0] == 'lessons':
        candidate = parts[2]
        if candidate in MATERIAL_TYPE_KEYS:
            return candidate
    return 'other'


def token_hash(token: str) -> str:
    return hashlib.sha256(token.encode('utf-8')).hexdigest()


def mail_configured() -> bool:
    return all([cfg.MAIL_SERVER, cfg.MAIL_USERNAME, cfg.MAIL_PASSWORD, cfg.MAIL_DEFAULT_SENDER])


def send_email(to_email: str, subject: str, text_body: str):
    if not mail_configured():
        raise RuntimeError('Mail is not configured.')
    message = EmailMessage()
    message['Subject'] = subject
    message['From'] = cfg.MAIL_DEFAULT_SENDER
    message['To'] = to_email
    message.set_content(text_body)

    if cfg.MAIL_USE_SSL:
        with smtplib.SMTP_SSL(cfg.MAIL_SERVER, cfg.MAIL_PORT, timeout=20) as smtp:
            smtp.login(cfg.MAIL_USERNAME, cfg.MAIL_PASSWORD)
            smtp.send_message(message)
        return

    with smtplib.SMTP(cfg.MAIL_SERVER, cfg.MAIL_PORT, timeout=20) as smtp:
        if cfg.MAIL_USE_TLS:
            smtp.starttls()
        smtp.login(cfg.MAIL_USERNAME, cfg.MAIL_PASSWORD)
        smtp.send_message(message)


def create_password_reset(db: Session, student: Student) -> str:
    token = secrets.token_urlsafe(32)
    reset = PasswordResetToken(
        student_id=student.id,
        token_hash=token_hash(token),
        expires_at=datetime.utcnow() + timedelta(minutes=PASSWORD_RESET_TOKEN_MINUTES),
    )
    db.add(reset)
    return token


def course_is_completed(db: Session, student_id: int, course: Course):
    lessons = db.query(Lesson.id).filter_by(course_id=course.id).all()
    lesson_ids = [row[0] for row in lessons]
    if not lesson_ids:
        return False
    completed = db.query(LessonProgress).filter(
        LessonProgress.student_id == student_id,
        LessonProgress.lesson_id.in_(lesson_ids),
        LessonProgress.is_completed.is_(True),
    ).count()
    return completed == len(lesson_ids) and course_quiz_average_passed(db, student_id, course.id)


def certificate_title(expertise_area: str, certificate_level: int):
    if certificate_level == MASTER_CERTIFICATE_LEVEL:
        return f'{expertise_area} Master Certificate'
    return f'{expertise_area} Level {certificate_level} Certificate'


def award_certificate(db: Session, student_id: int, expertise_area: str, certificate_level: int, hours_completed: int, source_course_id: int = None):
    existing = db.query(CertificateAward).filter_by(
        student_id=student_id,
        expertise_area=expertise_area,
        certificate_level=certificate_level,
    ).first()
    if existing:
        return existing
    award = CertificateAward(
        student_id=student_id,
        expertise_area=expertise_area,
        certificate_level=certificate_level,
        title=certificate_title(expertise_area, certificate_level),
        hours_completed=hours_completed,
        source_course_id=source_course_id,
        verification_code=uuid.uuid4().hex,
        issued_at=datetime.utcnow(),
    )
    db.add(award)
    return award


def evaluate_certificates(db: Session, student_id: int, course: Course):
    expertise_area = course.expertise_area or (course.program.name if course.program else None)
    certificate_level = course.certificate_level or 0
    if not expertise_area or certificate_level not in (1, 2, 3):
        return []
    if not course_is_completed(db, student_id, course):
        return []

    completed_courses = db.query(Course).join(Lesson, Lesson.course_id == Course.id).join(
        LessonProgress, LessonProgress.lesson_id == Lesson.id
    ).filter(
        Course.expertise_area == expertise_area,
        Course.certificate_level == certificate_level,
        LessonProgress.student_id == student_id,
        LessonProgress.is_completed.is_(True),
    ).distinct().all()
    hours_completed = 0
    for completed_course in completed_courses:
        if course_is_completed(db, student_id, completed_course):
            hours_completed += completed_course.learning_hours or 0
    if hours_completed < CERTIFICATE_LEVEL_HOURS:
        return []

    awards = [award_certificate(db, student_id, expertise_area, certificate_level, hours_completed, course.id)]
    completed_level_awards = db.query(CertificateAward).filter(
        CertificateAward.student_id == student_id,
        CertificateAward.expertise_area == expertise_area,
        CertificateAward.certificate_level.in_([1, 2, 3]),
    ).count()
    if completed_level_awards == 3:
        awards.append(award_certificate(db, student_id, expertise_area, MASTER_CERTIFICATE_LEVEL, CERTIFICATE_LEVEL_HOURS * 3, course.id))
    return awards


@app.get('/', response_class=HTMLResponse)
def home(request: Request, db: Session = Depends(get_db)):
    featured = db.query(Course).filter_by(is_published=True, is_featured=True).order_by(Course.created_at.desc()).limit(6).all()
    courses = db.query(Course).filter_by(is_published=True).order_by(Course.created_at.desc()).limit(12).all()
    # Courses per certificate level, grouped by expertise area (for the Level modals)
    level_courses = {}
    for level in (1, 2, 3):
        rows = db.query(Course).filter_by(is_published=True, certificate_level=level).order_by(
            Course.expertise_area, Course.title).all()
        groups = {}
        for course in rows:
            groups.setdefault(course.expertise_area or 'Other', []).append(course)
        level_courses[level] = groups
    return template(request, 'home.html', db, {'featured': featured, 'courses': courses,
                                               'categories': categories(db), 'level_courses': level_courses})


@app.get('/courses', response_class=HTMLResponse)
def catalog(request: Request, q: str = '', level: str = '', expertise: str = '', price: str = '', duration: str = '', db: Session = Depends(get_db)):
    query = db.query(Course).filter_by(is_published=True)
    if level:
        level_number_match = re.search(r'\d+', level)
        if level_number_match:
            query = query.filter(or_(Course.level == level, Course.certificate_level == int(level_number_match.group(0))))
        else:
            query = query.filter(Course.level == level)
    if expertise:
        query = query.filter(Course.expertise_area == expertise)
    if price == 'free':
        query = query.filter(or_(Course.price_cents <= 0, Course.allow_free_enrollment.is_(True)))
    elif price == 'paid':
        query = query.filter(Course.price_cents > 0, Course.allow_free_enrollment.is_(False))
    elif price == 'under-100':
        query = query.filter(Course.price_cents > 0, Course.price_cents < 10000)
    elif price == '100-500':
        query = query.filter(Course.price_cents >= 10000, Course.price_cents <= 50000)
    elif price == '500-plus':
        query = query.filter(Course.price_cents > 50000)
    if duration == 'short':
        query = query.filter(Course.learning_hours > 0, Course.learning_hours <= 5)
    elif duration == 'standard':
        query = query.filter(or_(Course.learning_hours.is_(None), Course.learning_hours == 0, and_(Course.learning_hours > 5, Course.learning_hours <= 15)))
    elif duration == 'extended':
        query = query.filter(Course.learning_hours > 15)
    if q:
        query = query.filter(or_(Course.title.ilike(f'%{q}%'), Course.description.ilike(f'%{q}%'), Course.sales_copy.ilike(f'%{q}%')))
    courses = query.order_by(Course.is_featured.desc(), Course.created_at.desc()).all()
    level_values = {row[0] for row in db.query(Course.level).filter(Course.is_published.is_(True), Course.level.isnot(None)).distinct().all() if row[0]}
    level_values.update(
        f'Level {row[0]}' for row in db.query(Course.certificate_level).filter(
            Course.is_published.is_(True),
            Course.certificate_level.in_([1, 2, 3]),
        ).distinct().all()
    )
    levels = sorted(level_values)
    return template(request, 'catalog.html', db, {
        'courses': courses,
        'levels': levels,
        'selected_level': level,
        'selected_expertise': expertise,
        'selected_price': price,
        'selected_duration': duration,
        'q': q,
        'categories': categories(db),
    })


@app.get('/api/expertise-areas')
def api_expertise_areas():
    return {'areas': EXPERTISE_AREAS}


@app.get('/api/certification-pathway')
def api_certification_pathway():
    return {
        'level_hours': CERTIFICATE_LEVEL_HOURS,
        'total_hours': CERTIFICATE_LEVEL_HOURS * 3,
        'levels': CERTIFICATE_LEVELS,
        'master_certificate_level': MASTER_CERTIFICATE_LEVEL,
    }


@app.get('/api/avatar-options')
def api_avatar_options():
    return {
        'text_enabled': True,
        'audio_enabled': bool(cfg.OPENAI_API_KEY),
        'video_enabled': bool(cfg.DID_API_KEY and cfg.DID_SOURCE_IMAGE_URL),
        'did_source_image_url': cfg.DID_SOURCE_IMAGE_URL if cfg.DID_SOURCE_IMAGE_URL else '',
        'openai_voice': cfg.OPENAI_VOICE,
    }


@app.get('/courses/{identifier}', response_class=HTMLResponse)
def course_detail(identifier: str, request: Request, db: Session = Depends(get_db)):
    course = find_course(db, identifier)
    if not course:
        raise HTTPException(status_code=404)
    student = student_from_request(request, db)
    enrolled = bool(student and db.query(Enrollment).filter_by(student_id=student.id, course_id=course.id, is_active=True).first())
    lessons = db.query(Lesson).filter_by(course_id=course.id).order_by(Lesson.lesson_number).all()
    return template(request, 'course_detail.html', db, {'course': course, 'lessons': lessons, 'enrolled': enrolled})


@app.get('/register', response_class=HTMLResponse)
def register_page(request: Request, next: str = '/courses'):
    with next_db_session() as db:
        return template(request, 'register.html', db, {'next_page': next})


@app.post('/register')
def register(request: Request, full_name: str = Form(...), email: str = Form(...), password: str = Form(...),
             confirm_password: str = Form(...), next: str = Form('/courses'), db: Session = Depends(get_db)):
    email = email.strip().lower()
    if len(password) < 10 or password != confirm_password:
        return RedirectResponse(f'/register?next={next}', status_code=303)
    if db.query(Student).filter_by(email=email).first():
        return RedirectResponse(f'/login?next={next}', status_code=303)
    student = Student(username=username_from_email(db, email), full_name=full_name.strip(), email=email, password_hash=hash_password(password), is_active=True)
    db.add(student)
    db.commit()
    db.refresh(student)
    request.session['student_id'] = student.id
    return RedirectResponse(next, status_code=303)


@app.get('/login', response_class=HTMLResponse)
def login_page(request: Request, next: str = '/learn/dashboard'):
    with next_db_session() as db:
        return template(request, 'login.html', db, {'next_page': next})


@app.post('/login')
def login(request: Request, username: str = Form(...), password: str = Form(...), next: str = Form('/learn/dashboard'), db: Session = Depends(get_db)):
    student = db.query(Student).filter((Student.username == username.strip()) | (Student.email == username.strip().lower())).first()
    if not student or not verify_password(password, student.password_hash):
        return RedirectResponse(f'/login?next={next}&login=failed', status_code=303)
    request.session['student_id'] = student.id
    student.last_login = datetime.utcnow()
    db.commit()
    return RedirectResponse(next, status_code=303)


@app.get('/forgot-password', response_class=HTMLResponse)
def forgot_password_page(request: Request):
    with next_db_session() as db:
        return template(request, 'forgot_password.html', db, {})


@app.post('/forgot-password')
def forgot_password(request: Request, email: str = Form(...), db: Session = Depends(get_db)):
    student = db.query(Student).filter_by(email=email.strip().lower(), is_active=True).first()
    if student and mail_configured():
        token = create_password_reset(db, student)
        db.commit()
        reset_url = f"{cfg.PUBLIC_BASE_URL.rstrip('/')}/reset-password/{token}"
        try:
            send_email(
                student.email,
                'Reset your Hub Academy password',
                (
                    f"Hello {student.full_name},\n\n"
                    "Use the link below to reset your Hub Academy password. "
                    f"This link expires in {PASSWORD_RESET_TOKEN_MINUTES} minutes.\n\n"
                    f"{reset_url}\n\n"
                    "If you did not request this, you can ignore this email."
                ),
            )
        except Exception as exc:
            logger.exception('Password reset email failed: %s', exc)
    elif student:
        logger.error('Password reset requested but mail is not configured.')
    return template(request, 'forgot_password_sent.html', db, {})


@app.get('/reset-password/{token}', response_class=HTMLResponse)
def reset_password_page(token: str, request: Request):
    with next_db_session() as db:
        reset = db.query(PasswordResetToken).filter_by(token_hash=token_hash(token), used_at=None).first()
        valid = bool(reset and reset.expires_at > datetime.utcnow())
        return template(request, 'reset_password.html', db, {'token': token if valid else '', 'valid': valid})


@app.post('/reset-password/{token}')
def reset_password(token: str, request: Request, password: str = Form(...), confirm_password: str = Form(...), db: Session = Depends(get_db)):
    reset = db.query(PasswordResetToken).filter_by(token_hash=token_hash(token), used_at=None).first()
    if not reset or reset.expires_at <= datetime.utcnow():
        return template(request, 'reset_password.html', db, {'token': '', 'valid': False})
    if len(password) < 10 or password != confirm_password:
        return template(request, 'reset_password.html', db, {'token': token, 'valid': True, 'error': 'Passwords must match and be at least 10 characters.'})

    student = db.get(Student, reset.student_id)
    if not student or not student.is_active:
        return template(request, 'reset_password.html', db, {'token': '', 'valid': False})
    student.password_hash = hash_password(password)
    reset.used_at = datetime.utcnow()
    db.commit()
    return RedirectResponse('/login', status_code=303)


@app.get('/logout')
def logout(request: Request):
    request.session.clear()
    return RedirectResponse('/', status_code=303)


@app.post('/courses/{identifier}/checkout')
def checkout(identifier: str, request: Request, db: Session = Depends(get_db)):
    student = student_from_request(request, db)
    if not student:
        return RedirectResponse(f'/login?next=/courses/{identifier}', status_code=303)
    course = find_course(db, identifier)
    if not course:
        raise HTTPException(status_code=404)
    amount = course.price_cents or 0
    currency = (course.currency or cfg.DEFAULT_CURRENCY).lower()
    if amount <= 0 or course.allow_free_enrollment:
        db.add(Purchase(student_id=student.id, course_id=course.id, amount_cents=0, currency=currency.upper(), status='paid', provider='free', completed_at=datetime.utcnow()))
        enroll_student(db, student.id, course.id)
        db.commit()
        return RedirectResponse(f'/learn/course/{course.id}', status_code=303)
    if not cfg.STRIPE_SECRET_KEY:
        return RedirectResponse(f'/courses/{course_slug(course)}', status_code=303)
    import stripe
    stripe.api_key = cfg.STRIPE_SECRET_KEY
    session = stripe.checkout.Session.create(
        mode='payment',
        customer_email=student.email,
        line_items=[{'price_data': {'currency': currency, 'product_data': {'name': course.title}, 'unit_amount': amount}, 'quantity': 1}],
        metadata={'student_id': student.id, 'course_id': course.id},
        success_url=f"{cfg.PUBLIC_BASE_URL.rstrip()}/checkout/success?session_id={{CHECKOUT_SESSION_ID}}",
        cancel_url=f"{cfg.PUBLIC_BASE_URL.rstrip()}/courses/{course_slug(course)}",
    )
    db.add(Purchase(student_id=student.id, course_id=course.id, amount_cents=amount, currency=currency.upper(), status='pending', provider='stripe', provider_session_id=session.id))
    db.commit()
    return RedirectResponse(session.url, status_code=303)


@app.post('/subscribe/{plan}')
def subscribe(plan: str, request: Request):
    with next_db_session() as db:
        student = student_from_request(request, db)
    if not student:
        return RedirectResponse('/login?next=/#pricing', status_code=303)
    selected = {p['name'].lower(): p for p in subscription_plans()}.get(plan.lower())
    if not selected or not cfg.STRIPE_SECRET_KEY or not selected['stripe_price_id']:
        return RedirectResponse('/#pricing', status_code=303)
    import stripe
    stripe.api_key = cfg.STRIPE_SECRET_KEY
    session = stripe.checkout.Session.create(
        mode='subscription',
        customer_email=student.email,
        line_items=[{'price': selected['stripe_price_id'], 'quantity': 1}],
        metadata={'student_id': student.id, 'plan': selected['name']},
        success_url=f"{cfg.PUBLIC_BASE_URL.rstrip()}/learn/dashboard",
        cancel_url=f"{cfg.PUBLIC_BASE_URL.rstrip()}/#pricing",
    )
    return RedirectResponse(session.url, status_code=303)


@app.get('/checkout/success')
def checkout_success(session_id: str = '', request: Request = None, db: Session = Depends(get_db)):
    purchase = db.query(Purchase).filter_by(provider_session_id=session_id).first()
    if purchase and purchase.status == 'paid':
        return RedirectResponse(f'/learn/course/{purchase.course_id}', status_code=303)
    return RedirectResponse('/learn/dashboard', status_code=303)


@app.post('/stripe/webhook')
async def stripe_webhook(request: Request, db: Session = Depends(get_db)):
    if not cfg.STRIPE_WEBHOOK_SECRET or not cfg.STRIPE_SECRET_KEY:
        return JSONResponse({'error': 'Stripe webhook is not configured'}, status_code=400)
    import stripe
    stripe.api_key = cfg.STRIPE_SECRET_KEY
    payload = await request.body()
    signature = request.headers.get('Stripe-Signature', '')
    try:
        event = stripe.Webhook.construct_event(payload, signature, cfg.STRIPE_WEBHOOK_SECRET)
    except Exception:
        return JSONResponse({'error': 'Invalid webhook'}, status_code=400)
    if event['type'] == 'checkout.session.completed':
        session = event['data']['object']
        purchase = db.query(Purchase).filter_by(provider_session_id=session['id']).first()
        if purchase:
            purchase.status = 'paid'
            purchase.provider_payment_intent = session.get('payment_intent')
            purchase.completed_at = datetime.utcnow()
            enroll_student(db, purchase.student_id, purchase.course_id)
            db.commit()
    return {'received': True}


@app.post('/assistant')
async def assistant(request: Request, db: Session = Depends(get_db)):
    data = await request.json()
    question = (data.get('question') or '').strip()
    mode = (data.get('mode') or 'text').strip().lower()
    history = data.get('history') or []
    courses = db.query(Course).filter_by(is_published=True).order_by(Course.is_featured.desc(), Course.created_at.desc()).limit(20).all()
    course_lines = [f"- {c.title} ({c.expertise_area or c.level or 'self-paced'}, {certificate_badge(c)}, {course_price(c)})" for c in courses]
    fallback = 'I can help you choose an expertise area, pick Level 1, Level 2, or Level 3 courses, and explain the Master Certificate path.'
    if not question:
        return {'answer': fallback, 'mode': mode, 'audio_available': False, 'video_available': False}
    try:
        client = openai_client(db)
        system = (
            "You are Hub Academy's AI learning guide. Answer the visitor's question directly and helpfully. "
            "Use a hybrid approach: draw on your general knowledge of the subject, skills, and careers, AND ground "
            "any course, level, or price recommendations in the live catalog below. Never invent courses that are not listed. "
            "When it fits, connect the answer to a relevant Hub Academy course or certificate level. Be concise, practical, and friendly.\n"
            "Format every answer as a short one-line intro, then 2-5 concise bullet points, each starting with '- '. "
            "Use **bold** for key terms and course names. Keep it scannable; avoid long paragraphs.\n\n"
            "Certification model: each expertise area has Level 1, Level 2, and Level 3 certificates. "
            "Each level requires 15 hours across videos, practices, simulations, and evaluation. "
            "Completing all three levels grants a Master Certificate, which is verifiable via a public verification ID.\n\n"
            "Available courses:\n" + ("\n".join(course_lines) if course_lines else "(no published courses yet)")
        )
        messages = [{'role': 'system', 'content': system}]
        # conversation memory: replay recent prior turns from the client
        for turn in history[-8:]:
            role = turn.get('role') if isinstance(turn, dict) else None
            content = (turn.get('content') or '').strip() if isinstance(turn, dict) else ''
            if role in ('user', 'assistant') and content:
                messages.append({'role': role, 'content': content[:1500]})
        messages.append({'role': 'user', 'content': question})
        # prefer gpt-4o, fall back to configured/mini models if unavailable
        models = []
        for model in ['gpt-4o', cfg.OPENAI_MODEL, 'gpt-4o-mini']:
            if model and model not in models:
                models.append(model)
        response, last_error = None, None
        for model in models:
            try:
                response = client.chat.completions.create(
                    model=model, messages=messages, max_tokens=450, temperature=0.5,
                )
                break
            except Exception as exc:
                last_error = exc
        if response is None:
            raise last_error or RuntimeError('OpenAI chat completion failed.')
        return {
            'answer': response.choices[0].message.content,
            'mode': mode,
            'audio_available': True,
            'video_available': bool(cfg.DID_API_KEY and cfg.DID_SOURCE_IMAGE_URL),
        }
    except Exception as exc:
        logger.exception('assistant failed: %s', exc)
        return {'answer': fallback, 'mode': mode, 'audio_available': False, 'video_available': False}


@app.post('/learn/feedback')
async def exercise_feedback(request: Request, db: Session = Depends(get_db)):
    """Give a learner concise feedback on an interactive exercise submission.
    Used by Application HTML exercises via a 'Get Feedback' button."""
    data = await request.json()
    title = (data.get('title') or 'this exercise').strip()
    criteria = (data.get('criteria') or '').strip()
    content = (data.get('content') or '').strip()
    if len(content) < 5:
        return JSONResponse({'error': 'Fill in your answers first, then ask for feedback.'}, status_code=400)
    try:
        client = openai_client(db)
        system = (
            "You are an expert reviewer and coach for a business course on AI agents. "
            "Give concise, encouraging, specific feedback addressed directly to the learner as 'You'. "
            "Judge the submission against the criteria provided and what the learner actually wrote. "
            "Return ONLY valid JSON: {\"strengths\": [\"You ...\"], \"improvements\": [\"You ...\"]} "
            "with 2 to 4 short one-sentence bullets in each list."
        )
        user = f"Exercise: {title}\n\nWhat good looks like (criteria):\n{criteria or 'Be clear, specific, and practical.'}\n\nThe learner's submission:\n{content[:6000]}"
        models = []
        for model in ['gpt-4o', cfg.OPENAI_MODEL, 'gpt-4o-mini']:
            if model and model not in models:
                models.append(model)
        response, last_error = None, None
        for model in models:
            try:
                response = client.chat.completions.create(
                    model=model,
                    messages=[{'role': 'system', 'content': system}, {'role': 'user', 'content': user}],
                    max_tokens=500, temperature=0.4,
                )
                break
            except Exception as exc:
                last_error = exc
        if response is None:
            raise last_error or RuntimeError('OpenAI chat completion failed.')
        parsed = parse_json_response(response.choices[0].message.content, {})

        def clean(value):
            if isinstance(value, str):
                value = [value]
            if not isinstance(value, list):
                return []
            return [str(item).strip() for item in value if str(item).strip()][:4]

        return {
            'strengths': clean(parsed.get('strengths')) or ['You made a solid start on this exercise.'],
            'improvements': clean(parsed.get('improvements')) or ['You could add more detail and be more specific.'],
        }
    except Exception as exc:
        logger.exception('exercise feedback failed: %s', exc)
        return JSONResponse({'error': 'Feedback is unavailable right now. Please try again.'}, status_code=502)


@app.get('/learn/dashboard')
def learner_dashboard(request: Request, db: Session = Depends(get_db)):
    student = student_from_request(request, db)
    if not student:
        return RedirectResponse('/login?next=/learn/dashboard', status_code=303)
    enrollments = db.query(Enrollment).filter_by(student_id=student.id, is_active=True).order_by(Enrollment.enrolled_at.desc()).all()
    certificates = db.query(CertificateAward).filter_by(student_id=student.id).order_by(CertificateAward.issued_at.desc()).all()
    cards = []
    for enrollment in enrollments:
        total = db.query(Lesson).filter_by(course_id=enrollment.course_id).count()
        done = db.query(LessonProgress).join(Lesson, LessonProgress.lesson_id == Lesson.id).filter(
            LessonProgress.student_id == student.id,
            Lesson.course_id == enrollment.course_id,
            LessonProgress.is_completed.is_(True)
        ).count()
        cards.append({
            'enrollment': enrollment,
            'total': total,
            'done': done,
            'pct': int(done / total * 100) if total else 0,
            'quiz': course_quiz_stats(db, student.id, enrollment.course_id),
        })
    return template(request, 'learn/dashboard.html', db, {'student': student, 'cards': cards, 'certificates': certificates})


@app.get('/learn/course/{course_id}')
def learner_course(course_id: int, request: Request, db: Session = Depends(get_db)):
    student = student_from_request(request, db)
    if not student:
        return RedirectResponse(f'/login?next=/learn/course/{course_id}', status_code=303)
    enrollment = db.query(Enrollment).filter_by(student_id=student.id, course_id=course_id, is_active=True).first()
    if not enrollment:
        return RedirectResponse('/learn/dashboard', status_code=303)
    course = db.get(Course, course_id)
    journey = journey_for_course(db, course, student.id)
    lessons = [session['lesson'] for module in journey['modules'] for session in module['sessions']]
    return template(request, 'learn/course.html', db, {'student': student, 'course': course, 'lessons': lessons, 'journey': journey, 'progress': journey['progress']})


@app.get('/learn/lesson/{lesson_id}')
def learner_lesson(lesson_id: int, request: Request, db: Session = Depends(get_db)):
    student = student_from_request(request, db)
    if not student:
        return RedirectResponse(f'/login?next=/learn/lesson/{lesson_id}', status_code=303)
    lesson = db.get(Lesson, lesson_id)
    if not lesson:
        raise HTTPException(status_code=404)
    enrollment = db.query(Enrollment).filter_by(student_id=student.id, course_id=lesson.course_id, is_active=True).first()
    if not enrollment:
        return RedirectResponse('/learn/dashboard', status_code=303)
    progress = db.query(LessonProgress).filter_by(student_id=student.id, lesson_id=lesson.id).first()
    if not progress:
        progress = LessonProgress(student_id=student.id, lesson_id=lesson.id, started_at=datetime.utcnow())
        db.add(progress)
    progress.last_accessed_at = datetime.utcnow()
    db.commit()
    materials = db.query(LessonMaterial).filter_by(lesson_id=lesson.id).order_by(LessonMaterial.upload_order).all()
    lessons = db.query(Lesson).filter_by(course_id=lesson.course_id).order_by(Lesson.lesson_number).all()
    grouped_materials = group_lesson_materials(materials)
    current_module_number = next((index for index, item in enumerate(lessons, start=1) if item.id == lesson.id), lesson_session_number(lesson))
    module_groups = module_material_groups(materials, current_module_number)
    module_sessions = module_session_groups(materials, current_module_number, session_objective_map(db, lesson.course_id))
    module_extras = module_extra_materials(materials)
    module_nav = module_nav_for_lessons(db, student.id, lessons, lesson.id)
    quiz_attempts = db.query(QuizAttempt).join(Quiz, QuizAttempt.quiz_id == Quiz.id).filter(
        QuizAttempt.student_id == student.id,
        Quiz.lesson_id == lesson.id,
    ).order_by(QuizAttempt.attempted_at.desc()).all()
    previous_lesson = next_lesson = None
    for index, item in enumerate(lessons):
        if item.id == lesson.id:
            previous_lesson = lessons[index - 1] if index > 0 else None
            next_lesson = lessons[index + 1] if index + 1 < len(lessons) else None
            break
    previous_progress = db.query(LessonProgress).filter_by(student_id=student.id, lesson_id=previous_lesson.id).first() if previous_lesson else None
    if previous_lesson and not session_unlocked(previous_progress):
        return RedirectResponse(f'/learn/course/{lesson.course_id}', status_code=303)
    steps = step_state(progress)
    return template(request, 'learn/lesson.html', db, {
        'student': student,
        'lesson': lesson,
        'materials': materials,
        'grouped_materials': grouped_materials,
        'module_groups': module_groups,
        'module_sessions': module_sessions,
        'module_extras': module_extras,
        'module_nav': module_nav,
        'current_module_number': current_module_number,
        'progress': progress,
        'steps': steps,
        'quiz_attempts': quiz_attempts,
        'quiz_attempts_remaining': max(0, MAX_LESSON_QUIZ_ATTEMPTS - len(quiz_attempts)),
        'quiz_pass_score': QUIZ_PASS_SCORE,
        'previous_lesson': previous_lesson,
        'next_lesson': next_lesson,
    })


@app.post('/learn/lesson/{lesson_id}/ai/flashcards')
def learner_flashcards(lesson_id: int, request: Request, db: Session = Depends(get_db)):
    student = student_from_request(request, db)
    lesson = db.get(Lesson, lesson_id)
    if not student or not lesson or not db.query(Enrollment).filter_by(student_id=student.id, course_id=lesson.course_id, is_active=True).first():
        return JSONResponse({'error': 'Not allowed'}, status_code=403)
    materials = db.query(LessonMaterial).filter_by(lesson_id=lesson.id).order_by(LessonMaterial.upload_order).all()
    corpus = lesson_learning_corpus(lesson, materials)
    if not corpus:
        return JSONResponse({'error': 'No readable lesson content was found for flashcards. Upload HTML, PDF, text, syllabus, case, simulation, or toolkit materials with extractable text.'}, status_code=400)
    prompt = f"""Create 8 serious flashcards from the extracted learning content below.

Rules:
- Use only the learning ideas, frameworks, decisions, examples, and application points in the extracted content.
- Never ask about file names, document names, folders, upload structure, or platform mechanics.
- Do not create generic dictionary definitions unless the content requires the term.
- Each card must test useful recall or application for a student.
- Answers must be concrete and educational, not one-word answers.

Lesson context:
{lesson_context_text(lesson, materials)}

Extracted learning content:
{corpus}

Return only valid JSON:
[
  {{"question":"specific question from the content", "answer":"clear answer grounded in the content"}}
]"""
    try:
        response = openai_chat_completion(
            openai_client(db, lesson.course),
            [{'role': 'user', 'content': prompt}],
            max_tokens=1200,
            temperature=0.7,
        )
        cards = parse_json_response(response.choices[0].message.content, [])
        if not isinstance(cards, list) or not cards:
            return JSONResponse({'error': 'OpenAI returned no flashcards. Try again.'}, status_code=502)
        return {'ok': True, 'flashcards': cards[:8] if isinstance(cards, list) else []}
    except Exception as exc:
        return ai_error_response(exc, 'Could not generate flashcards.')


@app.post('/learn/lesson/{lesson_id}/ai/audio-summary')
def learner_audio_summary(lesson_id: int, request: Request, db: Session = Depends(get_db)):
    student = student_from_request(request, db)
    lesson = db.get(Lesson, lesson_id)
    if not student or not lesson or not db.query(Enrollment).filter_by(student_id=student.id, course_id=lesson.course_id, is_active=True).first():
        return JSONResponse({'error': 'Not allowed'}, status_code=403)
    materials = db.query(LessonMaterial).filter_by(lesson_id=lesson.id).order_by(LessonMaterial.upload_order).all()
    corpus = lesson_learning_corpus(lesson, materials)
    prompt = f"""Write a professional spoken audio summary for this lesson. Keep it under 350 words, practical, and specific.

Never summarize file names, folders, upload structure, or platform mechanics. Summarize the actual learning content.

Lesson context:
{lesson_context_text(lesson, materials)}

Extracted learning content:
{corpus or 'No readable material text was extracted. Use only the lesson title and description.'}"""
    try:
        client = openai_client(db, lesson.course)
        script_response = openai_chat_completion(
            client,
            [{'role': 'user', 'content': prompt}],
            max_tokens=800,
            temperature=0.6,
        )
        script = (script_response.choices[0].message.content or '').strip()
        if not script:
            return JSONResponse({'error': 'OpenAI returned no audio script. Try again.'}, status_code=502)
        audio_models = [cfg.OPENAI_AUDIO_MODEL]
        if cfg.OPENAI_AUDIO_MODEL != 'tts-1':
            audio_models.append('tts-1')
        speech = None
        last_audio_error = None
        for audio_model in audio_models:
            try:
                speech = client.audio.speech.create(
                    model=audio_model,
                    voice=cfg.OPENAI_VOICE,
                    input=script,
                )
                break
            except Exception as audio_exc:
                last_audio_error = audio_exc
        if speech is None:
            raise last_audio_error or RuntimeError('Audio generation failed.')
        return Response(content=speech.content, media_type='audio/mpeg')
    except Exception as exc:
        return ai_error_response(exc, 'Could not generate audio summary.')


@app.post('/learn/lesson/{lesson_id}/quiz/generate')
def learner_generate_quiz(lesson_id: int, request: Request, db: Session = Depends(get_db)):
    student = student_from_request(request, db)
    lesson = db.get(Lesson, lesson_id)
    if not student or not lesson or not db.query(Enrollment).filter_by(student_id=student.id, course_id=lesson.course_id, is_active=True).first():
        return JSONResponse({'error': 'Not allowed'}, status_code=403)
    attempt_count = db.query(QuizAttempt).join(Quiz, QuizAttempt.quiz_id == Quiz.id).filter(
        QuizAttempt.student_id == student.id,
        Quiz.lesson_id == lesson.id,
    ).count()
    if attempt_count >= MAX_LESSON_QUIZ_ATTEMPTS:
        return JSONResponse({'error': 'Maximum 3 quiz attempts reached for this lesson.'}, status_code=400)
    materials = db.query(LessonMaterial).filter_by(lesson_id=lesson.id).order_by(LessonMaterial.upload_order).all()
    previous_questions = []
    previous_quizzes = db.query(Quiz).filter_by(lesson_id=lesson.id, is_ai_generated=True).order_by(Quiz.created_at.desc()).limit(10).all()
    for quiz in previous_quizzes:
        for question in parse_json_response(quiz.questions_json, []):
            if isinstance(question, dict):
                previous_questions.append(question.get('question', ''))
    try:
        questions = ai_generate_quiz_questions(db, lesson, materials, previous_questions)
        if not questions:
            return JSONResponse({'error': 'Could not generate quiz questions.'}, status_code=502)
        quiz = Quiz(
            lesson_id=lesson.id,
            title=f'AI Quiz - {lesson.title}',
            questions_json=json.dumps(questions, ensure_ascii=False),
            is_ai_generated=True,
            language='en',
        )
        db.add(quiz)
        db.commit()
        db.refresh(quiz)
        public_questions = [{**q, 'correct_answer': None, 'explanation': None} for q in questions]
        return {'ok': True, 'quiz_id': quiz.id, 'questions': public_questions, 'attempts_remaining': MAX_LESSON_QUIZ_ATTEMPTS - attempt_count}
    except Exception as exc:
        return ai_error_response(exc, 'Could not generate quiz.')


@app.post('/learn/quiz/{quiz_id}/submit')
async def learner_submit_quiz(quiz_id: int, request: Request, db: Session = Depends(get_db)):
    student = student_from_request(request, db)
    quiz = db.get(Quiz, quiz_id)
    if not student or not quiz or not quiz.lesson:
        return JSONResponse({'error': 'Not allowed'}, status_code=403)
    if not db.query(Enrollment).filter_by(student_id=student.id, course_id=quiz.lesson.course_id, is_active=True).first():
        return JSONResponse({'error': 'Not allowed'}, status_code=403)
    data = await request.json()
    answers = data.get('answers') or {}
    keep_quiz = bool(data.get('keep_quiz', True))
    questions = parse_json_response(quiz.questions_json, [])
    total = len(questions) or 1
    correct = 0
    for index, question in enumerate(questions):
        expected = int(question.get('correct_answer', -1))
        try:
            actual = int(answers.get(str(index), answers.get(index, -999)))
        except Exception:
            actual = -999
        if actual == expected:
            correct += 1
    score = int(round(correct / total * 100))
    try:
        ai_feedback = ai_generate_feedback(db, quiz.lesson, score, questions, answers)
    except Exception as exc:
        logger.exception('Quiz feedback failed: %s', exc)
        ai_feedback = {
            'strengths': ['You completed the quiz.'],
            'improvements': ['You should review the explanations for the questions you missed.'],
            'feedback': 'You completed the quiz.',
            'recommendations': 'You should review the explanations for the questions you missed.',
        }
    attempt = QuizAttempt(
        student_id=student.id,
        quiz_id=quiz.id,
        score=score,
        answers_json=json.dumps(answers, ensure_ascii=False),
        questions_snapshot_json=json.dumps(questions, ensure_ascii=False),
        feedback=ai_feedback['feedback'],
        recommendations=ai_feedback['recommendations'],
    )
    db.add(attempt)
    progress = db.query(LessonProgress).filter_by(student_id=student.id, lesson_id=quiz.lesson_id).first()
    if not progress:
        progress = LessonProgress(student_id=student.id, lesson_id=quiz.lesson_id, started_at=datetime.utcnow())
        db.add(progress)
    progress.quiz_completed = score >= QUIZ_PASS_SCORE
    if score >= QUIZ_PASS_SCORE:
        progress.content_viewed = True
        progress.revise_viewed = True
        progress.is_completed = True
        progress.completed_at = progress.completed_at or datetime.utcnow()
        evaluate_certificates(db, student.id, quiz.lesson.course)
    if not keep_quiz:
        quiz.title = f'Discarded quiz - {quiz.lesson.title}'
    db.commit()
    return {
        'ok': True,
        'score': score,
        'passed': score >= QUIZ_PASS_SCORE,
        'strengths': ai_feedback.get('strengths', []),
        'improvements': ai_feedback.get('improvements', []),
        'feedback': ai_feedback['feedback'],
        'recommendations': ai_feedback['recommendations'],
        'questions': questions,
    }


# ============ MY LEARNING PROFILE (learning-capacity assessment) ============
import random
import learner_profiling as lp


def profile_public_tasks():
    """Capacity task list for the browser — options are shuffled so the correct
    answer isn't always in the same position, and the answer is never sent."""
    tasks = []
    for task in lp.CAPACITY_TASKS:
        item = {'id': task['id'], 'level': task['level'], 'difficulty': task['difficulty'],
                'type': task['type'], 'prompt': task['prompt'], 'hint': task.get('hint', '')}
        if task['type'] == 'mcq':
            options = task['options'][:]
            random.shuffle(options)
            item['options'] = options
        tasks.append(item)
    return tasks


def _capacity_correct_text(task):
    return task['options'][task['answer']]


CAPACITY_DIMENSIONS = ['comprehension', 'reasoning', 'decision', 'problem_solving', 'transfer']


def ai_generate_capacity_tasks(db):
    """AI-generate a fresh set of general learning-capacity tasks (never the same
    twice). Falls back to the built-in bank if the AI is unavailable."""
    prompt = (
        "Create 7 short questions that measure a person's general LEARNING CAPACITY (not tied to any subject). "
        "Cover: 2 comprehension, 2 reasoning (logic or decisions), 3 problem-solving/transfer (novel problems or "
        "applying an idea to a brand-new situation). Make 6 multiple-choice with exactly 4 options and one correct "
        "answer, and VARY which position is correct. Make the LAST one an OPEN question where the learner writes 2-3 "
        "sentences applying an idea to a new situation, and give a short grading rubric. Increase difficulty from 1 "
        "(easy) to 5 (hard) across the set.\n"
        "Return only valid JSON, an array of objects like:\n"
        '{"dimension":"comprehension","difficulty":1,"type":"mcq","question":"...","options":["..","..","..",".."],"correct_answer":0,"hint":"..."}\n'
        'and the last: {"dimension":"transfer","difficulty":5,"type":"open","question":"...","rubric":"what full/partial/low credit looks like","hint":"..."}\n'
        'dimension must be one of: comprehension, reasoning, decision, problem_solving, transfer.'
    )
    response = openai_chat_completion(openai_client(db), [{'role': 'user', 'content': prompt}],
                                      max_tokens=1900, temperature=0.7)
    data = parse_json_response(response.choices[0].message.content, [])
    tasks = []
    for index, raw in enumerate(data if isinstance(data, list) else []):
        if not isinstance(raw, dict):
            continue
        dimension = raw.get('dimension') if raw.get('dimension') in CAPACITY_DIMENSIONS else 'reasoning'
        try:
            difficulty = max(1, min(5, int(raw.get('difficulty', index + 1))))
        except (TypeError, ValueError):
            difficulty = min(5, index + 1)
        kind = 'open' if raw.get('type') == 'open' else 'mcq'
        item = {'id': f'cap{index + 1}', 'dimension': dimension, 'difficulty': difficulty, 'type': kind,
                'level': dimension.replace('_', ' ').title(), 'prompt': str(raw.get('question', '')).strip(),
                'hint': str(raw.get('hint', '') or '')}
        if kind == 'open':
            item['rubric'] = str(raw.get('rubric', '') or 'Full credit: correctly applies the idea to the new case.')
        else:
            options = [str(o) for o in (raw.get('options') or []) if str(o).strip()]
            if len(options) < 2:
                continue
            try:
                correct = max(0, min(len(options) - 1, int(raw.get('correct_answer', 0))))
            except (TypeError, ValueError):
                correct = 0
            item['options'] = options
            item['answer'] = correct
        tasks.append(item)
    return tasks


def ai_generate_topic_questions(db, topic, count=4):
    """Generate `count` topic-knowledge MCQs (easy→hard) to place the learner.
    Returns dicts with the correct answer index; caller shuffles for display."""
    prompt = (
        f'Create exactly {count} multiple-choice questions that test how much someone already knows about "{topic}".\n'
        f'Order them from easiest (question 1) to hardest (question {count}). Each question has exactly 4 options and '
        'one correct answer. Make the distractors plausible. Vary the position of the correct answer.\n\n'
        'Return only valid JSON:\n'
        '[{"question": "...", "options": ["..","..","..",".."], "correct_answer": 0, "difficulty": 1}]\n'
        'where difficulty increases from 1 to 5.'
    )
    response = openai_chat_completion(openai_client(db), [{'role': 'user', 'content': prompt}],
                                      max_tokens=1300, temperature=0.5)
    data = parse_json_response(response.choices[0].message.content, [])
    questions = []
    for index, raw in enumerate(data if isinstance(data, list) else []):
        if not isinstance(raw, dict):
            continue
        options = [str(o) for o in (raw.get('options') or []) if str(o).strip()]
        if len(options) < 2:
            continue
        try:
            correct = int(raw.get('correct_answer', 0))
        except (TypeError, ValueError):
            correct = 0
        correct = max(0, min(len(options) - 1, correct))
        try:
            difficulty = int(raw.get('difficulty', index + 1))
        except (TypeError, ValueError):
            difficulty = index + 1
        questions.append({'id': f'topic{index + 1}', 'question': str(raw.get('question', '')).strip(),
                          'options': options, 'answer': correct, 'difficulty': max(1, min(5, difficulty))})
    return questions


@app.post('/learn/profile/topic-questions')
async def learner_profile_topic_questions(request: Request, db: Session = Depends(get_db)):
    """Return topic-knowledge questions for the chosen topic (public). The
    correct answers are stashed in the session for grading, never sent to the browser."""
    data = await request.json()
    topic = (data.get('topic') or '').strip()[:200]
    if not topic:
        return {'questions': []}
    try:
        generated = ai_generate_topic_questions(db, topic)
    except Exception as exc:
        logger.info('Topic-question generation unavailable: %s', exc)
        generated = []
    stored, public = [], []
    for question in generated:
        correct_text = question['options'][question['answer']]
        shuffled = question['options'][:]
        random.shuffle(shuffled)
        stored.append({'id': question['id'], 'correct_text': correct_text, 'difficulty': question['difficulty']})
        public.append({'id': question['id'], 'type': 'mcq', 'level': 'Topic knowledge',
                       'difficulty': question['difficulty'], 'prompt': question['question'], 'options': shuffled})
    request.session['topic_quiz'] = stored

    # Also generate the general learning-capacity tasks fresh (never the same).
    try:
        capacity = ai_generate_capacity_tasks(db) or list(lp.CAPACITY_TASKS)
    except Exception as exc:
        logger.info('Capacity-task generation unavailable, using built-in bank: %s', exc)
        capacity = list(lp.CAPACITY_TASKS)
    cap_store, cap_public = [], []
    for task in capacity:
        meta = {'id': task['id'], 'dimension': task['dimension'], 'difficulty': task['difficulty'], 'type': task['type']}
        pub = {'id': task['id'], 'type': task['type'], 'level': task.get('level', ''),
               'difficulty': task['difficulty'], 'prompt': task['prompt'], 'hint': task.get('hint', '')}
        if task['type'] == 'mcq':
            meta['correct_text'] = task['options'][task['answer']][:160]
            options = task['options'][:]
            random.shuffle(options)
            pub['options'] = options
        else:
            meta['rubric'] = (task.get('rubric') or '')[:300]
            meta['prompt'] = (task.get('prompt') or '')[:200]
        cap_store.append(meta)
        cap_public.append(pub)
    request.session['capacity_tasks'] = cap_store
    return {'questions': public, 'capacity': cap_public, 'topic': topic}


def ai_grade_open_task(db, task, answer):
    """Grade an open 'apply it' answer 0..1 via OpenAI; degrade gracefully."""
    answer = (answer or '').strip()
    if not answer:
        return 0.0, 'No answer provided.'
    try:
        prompt = (
            "You are grading a short open response that tests whether a learner can APPLY an idea "
            "to a new case. Grade fairly and briefly.\n\n"
            f"Task: {task['prompt']}\n\nGrading rubric: {task['rubric']}\n\n"
            f"Learner's answer: {answer}\n\n"
            'Return only valid JSON: {"score": 0.0-1.0, "note": "one short sentence of feedback"}'
        )
        response = openai_chat_completion(openai_client(db), [{'role': 'user', 'content': prompt}],
                                          max_tokens=200, temperature=0.0)
        data = parse_json_response(response.choices[0].message.content, {})
        score = max(0.0, min(1.0, float(data.get('score', 0))))
        return score, (data.get('note') or '')[:200]
    except Exception as exc:
        logger.info('Open-task AI grading unavailable, using heuristic: %s', exc)
        words = len(answer.split())
        frac = 0.6 if words >= 20 else 0.4 if words >= 8 else 0.2
        return frac, 'Graded automatically (AI grader was unavailable).'


def _profile_result_ctx(db, topic, scores, strategy, band, topic_band):
    published = db.query(Course).filter_by(is_published=True).all()
    plan = lp.build_learning_plan(published, topic, topic_band or 'Beginner')
    return {'has_result': True, 'topic': topic, 'scores': scores, 'strategy': strategy,
            'band': band, 'topic_band': topic_band or 'Beginner', 'plan': plan}


@app.get('/learn/profile', response_class=HTMLResponse)
def learner_profile(request: Request, db: Session = Depends(get_db)):
    """Public. Guests can take the assessment; results live in the session until
    they log in, at which point we save the profile to their account."""
    student = student_from_request(request, db)
    guest_result = request.session.get('learning_profile')
    if student:
        profile = db.query(LearnerProfile).filter_by(student_id=student.id).first()
        # a guest who just finished the assessment and then logged in — keep their result
        if not profile and guest_result:
            profile = LearnerProfile(
                student_id=student.id, topic=guest_result.get('topic'),
                scores_json=json.dumps(guest_result.get('scores', {}), ensure_ascii=False),
                strategy_json=json.dumps(guest_result.get('strategy', {}), ensure_ascii=False),
                strategy_key=(guest_result.get('strategy') or {}).get('key'),
                level_band=guest_result.get('band'), topic_band=guest_result.get('topic_band'))
            db.add(profile)
            db.commit()
        request.session.pop('learning_profile', None)
        if profile:
            ctx = _profile_result_ctx(db, profile.topic, parse_json_response(profile.scores_json, {}),
                                      parse_json_response(profile.strategy_json, {}), profile.level_band,
                                      profile.topic_band)
            return template(request, 'learn/profile.html', db, {'student': student, 'is_guest': False, **ctx})
        return template(request, 'learn/profile.html', db, {'student': student, 'is_guest': False, 'has_result': False})
    # anonymous visitor
    if guest_result:
        ctx = _profile_result_ctx(db, guest_result.get('topic'), guest_result.get('scores', {}),
                                  guest_result.get('strategy', {}), guest_result.get('band'),
                                  guest_result.get('topic_band'))
        return template(request, 'learn/profile.html', db, {'student': None, 'is_guest': True, **ctx})
    return template(request, 'learn/profile.html', db, {'student': None, 'is_guest': True, 'has_result': False})


@app.get('/learn/profile/assessment', response_class=HTMLResponse)
def learner_profile_assessment(request: Request, db: Session = Depends(get_db)):
    # Public — no login required to take the self-evaluation.
    student = student_from_request(request, db)
    return template(request, 'learn/profile_assessment.html', db,
                    {'student': student, 'tasks': profile_public_tasks()})


@app.post('/learn/profile/submit')
async def learner_profile_submit(request: Request, db: Session = Depends(get_db)):
    student = student_from_request(request, db)
    data = await request.json()
    topic = (data.get('topic') or '').strip()[:200]

    # ---- Capacity tasks: AI-generated, graded from the session where answers live ----
    cap_meta = request.session.get('capacity_tasks')
    if not cap_meta:
        cap_meta = []
        for t in lp.CAPACITY_TASKS:
            meta = {'id': t['id'], 'dimension': t['dimension'], 'difficulty': t['difficulty'], 'type': t['type']}
            if t['type'] == 'mcq':
                meta['correct_text'] = t['options'][t['answer']]
            else:
                meta['rubric'] = t.get('rubric', '')
                meta['prompt'] = t.get('prompt', '')
            cap_meta.append(meta)
    by_id = {t['id']: t for t in cap_meta}
    responses, open_grades = [], {}
    for raw in (data.get('responses') or []):
        task = by_id.get(raw.get('id'))
        if not task:
            continue
        confidence = raw.get('confidence')
        entry = {'id': task['id'], 'time_ms': raw.get('time_ms'),
                 'hint_used': bool(raw.get('hint_used')),
                 'confidence': confidence if confidence in (1, 2, 3) else None}
        answer_text = (raw.get('answer') or '')
        if task['type'] == 'mcq':
            entry['answer'] = answer_text
            entry['correct'] = (answer_text.strip() == (task.get('correct_text') or '').strip())
        else:
            answer_text = answer_text[:2000]
            fraction, note = ai_grade_open_task(db, {'prompt': task.get('prompt', ''), 'rubric': task.get('rubric', '')}, answer_text)
            open_grades[task['id']] = fraction
            entry['answer_text'] = answer_text
            entry['open_score'] = fraction
            entry['open_note'] = note
        responses.append(entry)

    scores = lp.compute_scores(responses, open_grades, tasks=cap_meta)
    band = lp.level_band(scores)
    key, strategy, rationale = lp.recommend_strategy(scores)
    strategy_payload = {**strategy, 'key': key, 'rationale': rationale}

    # ---- Topic-knowledge placement (graded against session-stored answers) ----
    topic_key = {q['id']: q for q in (request.session.get('topic_quiz') or [])}
    correct_w = total_w = 0.0
    for raw in (data.get('topic_responses') or []):
        question = topic_key.get(raw.get('id'))
        if not question:
            continue
        weight = question.get('difficulty', 1)
        total_w += weight
        if (raw.get('answer') or '').strip() == (question.get('correct_text') or '').strip():
            correct_w += weight
    topic_band = lp.topic_band_from(correct_w / total_w) if total_w else 'Beginner'
    request.session.pop('topic_quiz', None)
    request.session.pop('capacity_tasks', None)

    if student:
        profile = db.query(LearnerProfile).filter_by(student_id=student.id).first()
        if not profile:
            profile = LearnerProfile(student_id=student.id)
            db.add(profile)
        profile.topic = topic
        profile.scores_json = json.dumps(scores, ensure_ascii=False)
        profile.responses_json = json.dumps(responses, ensure_ascii=False)
        profile.strategy_key = key
        profile.strategy_json = json.dumps(strategy_payload, ensure_ascii=False)
        profile.level_band = band
        profile.topic_band = topic_band
        db.commit()
        request.session.pop('learning_profile', None)
    else:
        # guest — hold the result in the session so it survives to the results page
        # and can be saved once they create an account (raw responses omitted to stay small)
        request.session['learning_profile'] = {
            'topic': topic, 'scores': scores, 'strategy': strategy_payload,
            'band': band, 'topic_band': topic_band,
        }
    return {'ok': True, 'redirect': '/learn/profile'}


# ==================== AI COURSE FACTORY ====================
FACTORY_PROGRAM_NAME = 'AI Course Factory'
LEVEL_TO_CERT = {'Beginner': 1, 'Intermediate': 2, 'Advanced': 3}
FACTORY_LESSON_COUNT = 4
DEFAULT_FACTORY_RULES = (
    "Each lesson must:\n"
    "- Open with a one-sentence summary of what the learner will be able to do.\n"
    "- Teach 5 to 7 substantial points, each explained in 2-4 sentences of plain professional language (not one-liners).\n"
    "- Include at least one concrete, realistic worked example or mini case.\n"
    "- Include at least TWO data visualizations drawn as inline SVG (for example a labelled donut/pie showing a "
    "breakdown and a horizontal bar chart comparing options), each wrapped in a <figure> with a <figcaption> caption.\n"
    "- Include at least one comparison <table> and one step-by-step framework/process diagram.\n"
    "- End with a slide of 3 key takeaways.\n"
    "- Match the stated level (Beginner / Intermediate / Advanced).\n"
    "- Be accurate and self-contained. Figures shown in charts must be clearly illustrative examples, never invented "
    "precise statistics or fake citations. No medical, legal or financial advice."
)


def factory_get_setting(db, key, default=''):
    row = db.query(Settings).filter_by(key=key).first()
    return row.value if row and row.value is not None else default


def factory_set_setting(db, key, value):
    row = db.query(Settings).filter_by(key=key).first()
    if not row:
        row = Settings(key=key)
        db.add(row)
    row.value = value


def factory_rules(db):
    return factory_get_setting(db, 'ai_factory_rules', DEFAULT_FACTORY_RULES) or DEFAULT_FACTORY_RULES


def factory_require_payment(db):
    return factory_get_setting(db, 'ai_factory_require_payment', '0') == '1'


def get_or_create_factory_program(db):
    program = db.query(Program).filter_by(name=FACTORY_PROGRAM_NAME).first()
    if not program:
        program = Program(name=FACTORY_PROGRAM_NAME, description='Personalized courses generated by the AI Course Factory.')
        db.add(program)
        db.flush()
    return program


def normalize_topic(topic):
    return re.sub(r'\s+', ' ', (topic or '').strip().lower())[:200]


def learner_profile_snapshot(request, db, student):
    """The learner's saved profile (scores + strategy + topic band) from their
    account or, for guests, from the session — whichever exists."""
    if student:
        profile = db.query(LearnerProfile).filter_by(student_id=student.id).first()
        if profile:
            return {'scores': parse_json_response(profile.scores_json, {}),
                    'strategy': parse_json_response(profile.strategy_json, {}),
                    'topic_band': profile.topic_band}
    guest = request.session.get('learning_profile')
    if guest:
        return {'scores': guest.get('scores', {}), 'strategy': guest.get('strategy', {}),
                'topic_band': guest.get('topic_band')}
    return None


def learner_brief(profile, level):
    """A plain-English brief the AI generates the course against, so content is
    tuned to the learner's cognitive scores and learning strategy."""
    if not profile:
        return f'Target level: {level}. No detailed learner profile is available; use clear, standard teaching.'
    scores = profile.get('scores', {})
    strategy = profile.get('strategy', {})
    knobs = ', '.join(strategy.get('knobs', []) or [])
    return (
        f"Target level: {level}. "
        f"Learning strategy: {strategy.get('name', 'balanced')} — {strategy.get('tagline', '')} "
        f"Preferred approach: {knobs or 'balanced practice'}. "
        f"Cognitive scores (0-100): knowledge {scores.get('knowledge', '?')}, reasoning {scores.get('reasoning', '?')}, "
        f"application {scores.get('application', '?')}, learning speed {scores.get('speed', '?')}. "
        "Adapt teaching to this: if application is low, add more worked examples and hands-on practice; "
        "if learning speed is high, be concise and add stretch challenges; if low, go step-by-step with more scaffolding; "
        "if reasoning is high, include deeper 'why' explanations. Always honour the preferred approach above."
    )


def strip_code_fence(text):
    cleaned = (text or '').strip()
    if cleaned.startswith('```'):
        cleaned = cleaned.split('\n', 1)[1] if '\n' in cleaned else ''
        if cleaned.endswith('```'):
            cleaned = cleaned[:-3]
    return cleaned.strip()


def ai_generate_course_blueprint(db, topic, level, brief=''):
    prompt = (
        f'Design a concise professional micro-course on "{topic}" for a {level}-level learner.\n'
        f'Tailor the course to this specific learner: {brief}\n'
        f'Return exactly {FACTORY_LESSON_COUNT} lessons that build in order.\n'
        'Return only valid JSON:\n'
        '{"title": "...", "description": "one sentence", "lessons": ['
        '{"title": "...", "objectives": ["...", "..."]}]}'
    )
    response = openai_chat_completion(openai_client(db), [{'role': 'user', 'content': prompt}],
                                      max_tokens=1000, temperature=0.5)
    data = parse_json_response(response.choices[0].message.content, {})
    lessons = data.get('lessons') if isinstance(data, dict) else None
    if not lessons:
        lessons = [{'title': f'{topic} — Part {i + 1}', 'objectives': []} for i in range(FACTORY_LESSON_COUNT)]
    return {
        'title': (data.get('title') if isinstance(data, dict) else None) or f'{topic} ({level})',
        'description': (data.get('description') if isinstance(data, dict) else None) or f'A personalized {level} course on {topic}.',
        'lessons': lessons[:FACTORY_LESSON_COUNT],
    }


LESSON_VISUAL_GUIDE = (
    'BUILD THESE VISUALS with real, topic-specific labels (adapt the examples — keep them valid):\n'
    '1) A donut/pie as inline SVG (segments must sum to 100). Copy and adapt this exact structure:\n'
    '<figure><svg viewBox="0 0 42 42" width="190" height="190" role="img" aria-label="breakdown">'
    '<circle cx="21" cy="21" r="15.9155" fill="none" stroke="#eef2fb" stroke-width="6"></circle>'
    '<circle cx="21" cy="21" r="15.9155" fill="none" stroke="#2563eb" stroke-width="6" stroke-dasharray="55 45" stroke-dashoffset="25"></circle>'
    '<circle cx="21" cy="21" r="15.9155" fill="none" stroke="#7c3aed" stroke-width="6" stroke-dasharray="30 70" stroke-dashoffset="-30"></circle>'
    '<circle cx="21" cy="21" r="15.9155" fill="none" stroke="#06b6d4" stroke-width="6" stroke-dasharray="15 85" stroke-dashoffset="-60"></circle>'
    '</svg><figcaption>Illustrative breakdown — Blue 55%, Purple 30%, Teal 15%</figcaption></figure>\n'
    '   (dashoffset for each segment = 25 minus the running total of the previous segments; label each colour in the caption.)\n'
    '2) A horizontal bar chart as inline SVG. Copy and adapt (one <text>+<rect> pair per bar; scale widths to your values):\n'
    '<figure><svg viewBox="0 0 340 160" width="100%" role="img" aria-label="comparison"><g font-size="11" fill="#33425c">'
    '<text x="0" y="24">Label A</text><rect x="96" y="13" width="210" height="15" rx="7" fill="#2563eb"></rect>'
    '<text x="0" y="58">Label B</text><rect x="96" y="47" width="150" height="15" rx="7" fill="#7c3aed"></rect>'
    '<text x="0" y="92">Label C</text><rect x="96" y="81" width="90" height="15" rx="7" fill="#06b6d4"></rect>'
    '</g></svg><figcaption>Illustrative comparison</figcaption></figure>\n'
    '3) A comparison table: <table><thead><tr><th>…</th><th>…</th></tr></thead><tbody><tr><td>…</td><td>…</td></tr></tbody></table>\n'
    '4) A step-by-step framework as: <div class="flow"><div class="flow-step"><b>1. Step</b><span>what happens</span></div> …3-5 steps… </div>\n'
    '5) Stat cards for memorable numbers/principles: <div class="stat-grid"><div class="stat"><b>3x</b><span>label</span></div> …2-4 cards… </div>\n'
    '6) A practical tip as: <div class="callout">One sharp, actionable tip.</div>\n'
)


def ai_generate_lesson_html(db, topic, level, lesson_title, objectives, rules, brief='', fix_notes=''):
    fix = f'\nThe previous attempt was rejected for: {fix_notes}. Fix these issues.' if fix_notes else ''
    prompt = (
        f'Write the full, richly-illustrated content for a {level}-level lesson titled "{lesson_title}" in a course on "{topic}".\n'
        f'Tailor the teaching to this specific learner: {brief}\n'
        f'Learning objectives: {json.dumps(objectives, ensure_ascii=False)}\n\n'
        f'Follow these rules exactly:\n{rules}\n{fix}\n\n'
        'Make it genuinely engaging and substantial: open with a one-line hook, use a vivid real-world analogy, keep '
        'paragraphs short and punchy but teach each point in depth, and include a concrete worked example the learner '
        'can picture. Produce 6 to 8 slides.\n\n'
        'Structure every slide as <section class="slide"> with a short <h2> title, a lead <p>, then the teaching '
        '(short paragraphs and <ul>/<li>). Use <blockquote> for the worked example or analogy. Distribute the required '
        'visuals across the slides so most slides carry a chart, table, diagram, stat cards or a callout — not walls of text. '
        'End with <section class="slide"><h2>Key takeaways</h2><ul>…3 items…</ul></section>.\n\n'
        f'{LESSON_VISUAL_GUIDE}\n'
        'Output clean semantic HTML only: no markdown, no code fences, no <html>/<head>/<body>, and NO <script> or '
        '<style> tags (inline style="" attributes and inline <svg> are allowed and encouraged). Every <svg> must be '
        'valid and self-contained.'
    )
    response = openai_chat_completion(openai_client(db), [{'role': 'user', 'content': prompt}],
                                      max_tokens=4200, temperature=0.6)
    html = strip_code_fence(response.choices[0].message.content)
    return re.sub(r'(?is)<(script|style)\b.*?</\1>', '', html)


def ai_review_lesson(db, lesson_title, html, rules):
    prompt = (
        'You are a strict course-quality reviewer. Check the lesson HTML against the rules and reply honestly.\n\n'
        f'Rules:\n{rules}\n\nLesson title: {lesson_title}\n\nLesson HTML:\n{html[:9000]}\n\n'
        'Return only valid JSON: {"pass": true/false, "issues": ["short issue", "..."]}'
    )
    try:
        response = openai_chat_completion(openai_client(db), [{'role': 'user', 'content': prompt}],
                                          max_tokens=300, temperature=0.0)
        data = parse_json_response(response.choices[0].message.content, {})
        return bool(data.get('pass', True)), (data.get('issues') or [])
    except Exception as exc:
        logger.info('Lesson review unavailable, accepting: %s', exc)
        return True, []


def factory_build_lesson(db, course, lesson, rules):
    """Generate + review one lesson, with up to 3 self-correcting attempts."""
    topic = course.source_topic or course.title
    level = course.source_level or 'Beginner'
    brief = course.generation_brief or ''
    objectives = parse_json_response(lesson.description, []) if lesson.description else []
    fix_notes = ''
    for _attempt in range(3):
        html = ai_generate_lesson_html(db, topic, level, lesson.title, objectives, rules, brief, fix_notes)
        passed, issues = ai_review_lesson(db, lesson.title, html, rules)
        if passed:
            lesson.content_html = html
            lesson.generation_status = 'ready'
            lesson.review_notes = 'Approved by AI reviewer.'
            return True
        fix_notes = '; '.join(issues)[:400]
        lesson.content_html = html  # keep the latest draft
    lesson.generation_status = 'needs_review'
    lesson.review_notes = 'Auto-review flagged: ' + (fix_notes or 'quality issues')
    return False


@app.post('/learn/factory/build')
async def learner_factory_build(request: Request, db: Session = Depends(get_db)):
    data = await request.json()
    topic = (data.get('topic') or '').strip()[:200]
    level = data.get('level') if data.get('level') in LEVEL_TO_CERT else 'Beginner'
    if not topic:
        return JSONResponse({'error': 'Please choose a topic.'}, status_code=400)
    student = student_from_request(request, db)
    if factory_require_payment(db) and not student:
        return JSONResponse({'error': 'login', 'redirect': '/login?next=/learn/profile'}, status_code=402)

    # Tailor generation to the learner's saved profile (cognitive scores + strategy)
    profile = learner_profile_snapshot(request, db, student)
    strategy_key = ((profile or {}).get('strategy') or {}).get('key') or 'generic'
    brief = learner_brief(profile, level)

    # Reuse if a course for this topic+level+profile-type already exists (built once, shared)
    norm = normalize_topic(topic)
    existing = db.query(Course).filter_by(is_ai_generated=True, source_topic=norm,
                                          source_level=level, source_profile=strategy_key).first()
    if existing and existing.generation_status != 'failed':
        if student:
            enroll_student(db, student.id, existing.id)
            db.commit()
        return {'ok': True, 'redirect': f'/learn/factory/course/{existing.id}'}

    program = get_or_create_factory_program(db)
    try:
        blueprint = ai_generate_course_blueprint(db, topic, level, brief)
    except Exception as exc:
        return ai_error_response(exc, 'Could not start course generation.')

    course = Course(program_id=program.id, title=blueprint['title'][:200], description=blueprint['description'],
                    slug=slugify(blueprint['title'])[:200] or f'ai-{norm}', expertise_area=topic[:120],
                    certificate_level=LEVEL_TO_CERT[level], learning_hours=FACTORY_LESSON_COUNT * 2,
                    num_lessons=len(blueprint['lessons']), is_published=False, is_ai_generated=True,
                    generation_status='building', source_topic=norm, source_level=level,
                    source_profile=strategy_key, generation_brief=brief, created_at=datetime.utcnow())
    db.add(course)
    db.flush()
    if db.query(Course).filter(Course.slug == course.slug, Course.id != course.id).first():
        course.slug = f'{course.slug}-{course.id}'
    for index, item in enumerate(blueprint['lessons']):
        db.add(Lesson(course_id=course.id, lesson_number=index + 1, module_number=index + 1,
                      title=(item.get('title') or f'Lesson {index + 1}')[:200],
                      description=json.dumps(item.get('objectives') or [], ensure_ascii=False),
                      generation_status='queued'))
    db.flush()

    # Build lesson 1 immediately as a free preview so the learner can start now
    first = db.query(Lesson).filter_by(course_id=course.id).order_by(Lesson.lesson_number).first()
    if first:
        first.generation_status = 'building'
        try:
            factory_build_lesson(db, course, first, factory_rules(db))
        except Exception as exc:
            logger.exception('First lesson build failed: %s', exc)
            first.generation_status = 'failed'
            first.review_notes = str(exc)[:300]
    # Freemium model: only Lesson 1 is generated as a preview. The rest stay
    # locked until the learner registers & subscribes; the full course is then
    # delivered to their account (produced separately).
    db.query(Lesson).filter(Lesson.course_id == course.id, Lesson.generation_status == 'queued')\
        .update({'generation_status': 'locked'}, synchronize_session=False)
    course.generation_status = 'sample'
    if student:
        enroll_student(db, student.id, course.id)
    db.commit()
    return {'ok': True, 'redirect': f'/learn/factory/course/{course.id}'}


def factory_sample_images(course):
    """A few topical images for the sample page — reuses the app's course art."""
    images = []
    try:
        images.append(course_image(course))
    except Exception:
        pass
    images.extend(DEFAULT_COURSE_IMAGES)
    seen, out = set(), []
    for url in images:
        if url and url not in seen:
            seen.add(url)
            out.append(url)
    return out


# Short (~25s) topic-intro videos rendered offline and shipped under
# /static/videos/factory/<slug>/. Keyed by normalized topic so it is portable
# across environments (course ids differ per DB). Empty -> no video section.
FACTORY_INTRO_VIDEOS = {
    'negotiations for leaders': {
        'src': '/static/videos/factory/negotiations-for-leaders/intro.mp4',
        'poster': '/static/videos/factory/negotiations-for-leaders/intro.jpg',
        'dur': '0:28'},
}


def factory_intro_video(course):
    """The single cinematic intro video for this course, if we've produced one."""
    key = normalize_topic(course.source_topic or course.expertise_area or course.title or '')
    return FACTORY_INTRO_VIDEOS.get(key)


import secrets as _secrets


@app.post('/api/factory/deliver')
async def factory_deliver_course(request: Request, db: Session = Depends(get_db)):
    """The external generator app delivers a finished course into this platform.

    Auth: header 'X-Factory-Key' must equal FACTORY_API_KEY.
    Body JSON:
      {
        "student_email": "user@example.com",     # required — the login they paid with
        "student_name":  "Full Name",            # optional (used if we create the account)
        "title":         "AI Strategy Mastery",  # required
        "description":   "...",
        "topic":         "AI Strategy",
        "level":         "Intermediate",         # Beginner|Intermediate|Advanced
        "lessons": [ {"title": "...", "content_html": "<section>...</section>",
                       "objectives": ["..."]}, ... ]   # required, 1+
      }
    Creates the course + lessons, enrols the student, and returns the course URL.
    The student then logs in with that email and finds the course in My Learning.
    """
    if not cfg.FACTORY_API_KEY:
        return JSONResponse({'error': 'Delivery API is not configured (set FACTORY_API_KEY).'}, status_code=503)
    if request.headers.get('X-Factory-Key') != cfg.FACTORY_API_KEY:
        return JSONResponse({'error': 'Unauthorized'}, status_code=401)
    data = await request.json()
    email = (data.get('student_email') or '').strip().lower()
    title = (data.get('title') or '').strip()
    lessons_in = data.get('lessons') or []
    if not email or not title or not lessons_in:
        return JSONResponse({'error': 'student_email, title and at least one lesson are required.'}, status_code=400)

    level = data.get('level') if data.get('level') in LEVEL_TO_CERT else 'Beginner'
    topic = (data.get('topic') or title).strip()[:200]

    student = db.query(Student).filter_by(email=email).first()
    created_account, temp_password = False, None
    if not student:
        temp_password = _secrets.token_urlsafe(9)
        student = Student(username=username_from_email(db, email),
                          full_name=(data.get('student_name') or email.split('@')[0]).strip()[:120],
                          email=email, password_hash=hash_password(temp_password), is_active=True)
        db.add(student)
        db.flush()
        created_account = True

    program = get_or_create_factory_program(db)
    course = Course(program_id=program.id, title=title[:200],
                    description=(data.get('description') or '').strip()[:2000] or None,
                    slug=slugify(title)[:200] or f'delivered-{normalize_topic(topic)}',
                    expertise_area=topic[:120], certificate_level=LEVEL_TO_CERT[level],
                    num_lessons=len(lessons_in), learning_hours=len(lessons_in) * 2,
                    is_published=False, is_ai_generated=True, generation_status='delivered',
                    source_topic=normalize_topic(topic), source_level=level,
                    source_profile=(data.get('profile') or 'delivered'), created_at=datetime.utcnow())
    db.add(course)
    db.flush()
    if db.query(Course).filter(Course.slug == course.slug, Course.id != course.id).first():
        course.slug = f'{course.slug}-{course.id}'
    for index, item in enumerate(lessons_in):
        db.add(Lesson(course_id=course.id, lesson_number=index + 1, module_number=index + 1,
                      title=(item.get('title') or f'Lesson {index + 1}')[:200],
                      description=json.dumps(item.get('objectives') or [], ensure_ascii=False),
                      content_html=item.get('content_html') or '', generation_status='ready'))
    enroll_student(db, student.id, course.id)
    db.commit()
    result = {'ok': True, 'course_id': course.id, 'course_url': f'/learn/factory/course/{course.id}',
              'student_id': student.id, 'account_created': created_account}
    if temp_password:
        result['temporary_password'] = temp_password
    return result


@app.post('/learn/factory/course/{course_id}/generate-next')
def learner_factory_generate_next(course_id: int, request: Request, db: Session = Depends(get_db)):
    course = db.get(Course, course_id)
    if not course or not course.is_ai_generated:
        raise HTTPException(status_code=404)
    nxt = db.query(Lesson).filter_by(course_id=course.id, generation_status='queued').order_by(Lesson.lesson_number).first()
    if nxt:
        nxt.generation_status = 'building'
        db.commit()
        try:
            factory_build_lesson(db, course, nxt, factory_rules(db))
        except Exception as exc:
            logger.exception('Lesson build failed: %s', exc)
            nxt.generation_status = 'failed'
            nxt.review_notes = str(exc)[:300]
    remaining = db.query(Lesson).filter(Lesson.course_id == course.id,
                                        Lesson.generation_status.in_(['queued', 'building'])).count()
    if remaining == 0 and course.generation_status == 'building':
        course.generation_status = 'ready'
    db.commit()
    lessons = db.query(Lesson).filter_by(course_id=course.id).order_by(Lesson.lesson_number).all()
    return {
        'done': remaining == 0,
        'lessons': [{'id': l.id, 'title': l.title, 'status': l.generation_status} for l in lessons],
    }


@app.get('/learn/factory/course/{course_id}', response_class=HTMLResponse)
def learner_factory_course(course_id: int, request: Request, db: Session = Depends(get_db)):
    course = db.get(Course, course_id)
    if not course or not course.is_ai_generated:
        raise HTTPException(status_code=404)
    lessons = db.query(Lesson).filter_by(course_id=course.id).order_by(Lesson.lesson_number).all()
    return template(request, 'learn/factory_course.html', db,
                    {'student': student_from_request(request, db), 'course': course, 'lessons': lessons,
                     'images': factory_sample_images(course)})


@app.get('/learn/factory/lesson/{lesson_id}', response_class=HTMLResponse)
def learner_factory_lesson(lesson_id: int, request: Request, db: Session = Depends(get_db)):
    lesson = db.get(Lesson, lesson_id)
    if not lesson or not lesson.course or not lesson.course.is_ai_generated:
        raise HTTPException(status_code=404)
    lessons = db.query(Lesson).filter_by(course_id=lesson.course_id).order_by(Lesson.lesson_number).all()
    objectives = parse_json_response(lesson.description, []) if lesson.description else []
    return template(request, 'learn/factory_lesson.html', db,
                    {'student': student_from_request(request, db), 'course': lesson.course,
                     'lesson': lesson, 'lessons': lessons, 'objectives': objectives,
                     'images': factory_sample_images(lesson.course),
                     'intro_video': factory_intro_video(lesson.course)})


def factory_lesson_narration_text(lesson):
    """Plain, speakable text from the lesson's HTML for Emma to read."""
    text = re.sub(r'(?s)<[^>]+>', ' ', lesson.content_html or '')
    for entity, char in [('&amp;', '&'), ('&lt;', '<'), ('&gt;', '>'), ('&#39;', "'"),
                         ('&quot;', '"'), ('&nbsp;', ' ')]:
        text = text.replace(entity, char)
    return re.sub(r'\s+', ' ', text).strip()


# Emma = a natural neural voice (OpenAI TTS). "nova" is a warm, non-robotic female voice.
EMMA_TTS_VOICE = 'nova'


@app.get('/learn/factory/lesson/{lesson_id}/audio')
def learner_factory_lesson_audio(lesson_id: int, request: Request, db: Session = Depends(get_db)):
    """Emma narration for a factory lesson — generated once via OpenAI TTS and cached in R2."""
    lesson = db.get(Lesson, lesson_id)
    if not lesson or not lesson.course or not lesson.course.is_ai_generated:
        raise HTTPException(status_code=404)
    if lesson.audio_key:
        try:
            return RedirectResponse(presigned_download_url(lesson.audio_key, f'emma-lesson-{lesson.id}.mp3'), status_code=303)
        except Exception:
            pass
    text = factory_lesson_narration_text(lesson)
    if not text:
        raise HTTPException(status_code=404)
    try:
        speech = openai_client(db).audio.speech.create(model='tts-1', voice=EMMA_TTS_VOICE, input=text[:3900])
        audio_bytes = getattr(speech, 'content', None) or speech.read()
    except Exception as exc:
        logger.exception('Emma TTS failed: %s', exc)
        return JSONResponse({'error': 'Voice narration is unavailable right now.'}, status_code=502)
    if r2_enabled():
        import io
        key = object_key(f'emma-{lesson.id}.mp3', material_type='factory-audio')
        try:
            upload_fileobj(key, io.BytesIO(audio_bytes), 'audio/mpeg')
            lesson.audio_key = key
            db.commit()
            return RedirectResponse(presigned_download_url(key, f'emma-lesson-{lesson.id}.mp3'), status_code=303)
        except Exception as exc:
            logger.info('Emma audio R2 store failed, streaming instead: %s', exc)
    from fastapi.responses import Response
    return Response(content=audio_bytes, media_type='audio/mpeg')


@app.get('/admin/ai-factory', response_class=HTMLResponse)
def admin_ai_factory(request: Request, db: Session = Depends(get_db)):
    require_admin(request, db)
    courses = db.query(Course).filter_by(is_ai_generated=True).order_by(Course.created_at.desc()).all()
    rows = []
    for course in courses:
        total = db.query(Lesson).filter_by(course_id=course.id).count()
        ready = db.query(Lesson).filter_by(course_id=course.id, generation_status='ready').count()
        needs = db.query(Lesson).filter_by(course_id=course.id, generation_status='needs_review').count()
        rows.append({'course': course, 'total': total, 'ready': ready, 'needs': needs})
    return template(request, 'admin/ai_factory.html', db, {'rows': rows})


@app.get('/admin/ai-factory/rules', response_class=HTMLResponse)
def admin_ai_factory_rules(request: Request, db: Session = Depends(get_db)):
    require_admin(request, db)
    return template(request, 'admin/ai_factory_rules.html', db,
                    {'rules': factory_rules(db), 'require_payment': factory_require_payment(db),
                     'saved': request.query_params.get('saved')})


@app.post('/admin/ai-factory/rules')
def admin_ai_factory_rules_save(request: Request, rules: str = Form(''), require_payment: str = Form(''),
                                db: Session = Depends(get_db)):
    require_admin(request, db)
    factory_set_setting(db, 'ai_factory_rules', rules.strip() or DEFAULT_FACTORY_RULES)
    factory_set_setting(db, 'ai_factory_require_payment', '1' if require_payment else '0')
    db.commit()
    return RedirectResponse('/admin/ai-factory/rules?saved=1', status_code=303)


@app.post('/admin/ai-factory/course/{course_id}/delete')
def admin_ai_factory_delete(course_id: int, request: Request, db: Session = Depends(get_db)):
    require_admin(request, db)
    course = db.get(Course, course_id)
    if course and course.is_ai_generated:
        try:
            cascade_delete_course(db, course)
            db.commit()
        except Exception:
            db.rollback()
    return RedirectResponse('/admin/ai-factory', status_code=303)


@app.post('/admin/ai-factory/course/{course_id}/rebuild')
def admin_ai_factory_rebuild(course_id: int, request: Request, db: Session = Depends(get_db)):
    """Regenerate the already-built lessons of an existing AI course in place, so
    older courses pick up the current prompt/rules (rich charts, tables, etc.).
    Locked freemium lessons are left untouched; Emma audio is cleared to re-narrate."""
    require_admin(request, db)
    course = db.get(Course, course_id)
    if course and course.is_ai_generated:
        rules = factory_rules(db)
        lessons = db.query(Lesson).filter_by(course_id=course.id).order_by(Lesson.lesson_number).all()
        for lesson in lessons:
            if lesson.generation_status in ('locked', 'queued'):
                continue  # not yet unlocked/built — leave as-is
            lesson.audio_key = None  # force Emma to re-narrate the new text
            lesson.generation_status = 'building'
            db.commit()
            try:
                factory_build_lesson(db, course, lesson, rules)
            except Exception as exc:
                logger.exception('Rebuild failed for lesson %s: %s', lesson.id, exc)
                lesson.generation_status = 'failed'
                lesson.review_notes = str(exc)[:300]
            db.commit()
    return RedirectResponse('/admin/ai-factory', status_code=303)


@app.post('/learn/lesson/{lesson_id}/step/{step_number}/complete')
def learner_complete_step(lesson_id: int, step_number: int, request: Request, db: Session = Depends(get_db)):
    student = student_from_request(request, db)
    if not student:
        return JSONResponse({'error': 'Login required'}, status_code=401)
    lesson = db.get(Lesson, lesson_id)
    if not lesson:
        raise HTTPException(status_code=404)
    enrollment = db.query(Enrollment).filter_by(student_id=student.id, course_id=lesson.course_id, is_active=True).first()
    if not enrollment:
        return JSONResponse({'error': 'Enrollment required'}, status_code=403)
    progress = db.query(LessonProgress).filter_by(student_id=student.id, lesson_id=lesson.id).first()
    if not progress:
        progress = LessonProgress(student_id=student.id, lesson_id=lesson.id, started_at=datetime.utcnow())
        db.add(progress)
    if step_number == 1:
        progress.content_viewed = True
    elif step_number == 2:
        if not progress.content_viewed:
            return JSONResponse({'error': 'Complete step 1 first'}, status_code=400)
        progress.revise_viewed = True
    elif step_number == 3:
        return JSONResponse({'error': 'Complete the AI quiz with at least 60% to finish this module.'}, status_code=400)
    else:
        return JSONResponse({'error': 'Invalid step'}, status_code=400)
    if progress.content_viewed and progress.revise_viewed and progress.quiz_completed:
        progress.is_completed = True
        progress.completed_at = progress.completed_at or datetime.utcnow()
        course = db.get(Course, lesson.course_id)
        if course:
            evaluate_certificates(db, student.id, course)
    progress.last_accessed_at = datetime.utcnow()
    db.commit()
    return {
        'ok': True,
        'is_completed': bool(progress.is_completed),
        'steps': step_state(progress),
    }


@app.post('/learn/lesson/{lesson_id}/complete')
def learner_complete_lesson(lesson_id: int, request: Request, db: Session = Depends(get_db)):
    student = student_from_request(request, db)
    if not student:
        return RedirectResponse(f'/login?next=/learn/lesson/{lesson_id}', status_code=303)
    progress = db.query(LessonProgress).filter_by(student_id=student.id, lesson_id=lesson_id).first()
    if not progress:
        progress = LessonProgress(student_id=student.id, lesson_id=lesson_id)
        db.add(progress)
    progress.content_viewed = True
    progress.revise_viewed = True
    progress.quiz_completed = True
    progress.is_completed = True
    progress.completed_at = datetime.utcnow()
    lesson = db.get(Lesson, lesson_id)
    if lesson:
        course = db.get(Course, lesson.course_id)
        if course:
            evaluate_certificates(db, student.id, course)
    db.commit()
    return RedirectResponse(f'/learn/course/{lesson.course_id}', status_code=303)


@app.get('/admin')
@app.get('/admin/dashboard')
def admin_dashboard(request: Request, db: Session = Depends(get_db)):
    admin = require_admin(request, db)
    stats = {
        'programs': db.query(Program).count(),
        'courses': db.query(Course).count(),
        'students': db.query(Student).count(),
        'enrollments': db.query(Enrollment).count(),
    }
    recent_courses = db.query(Course).order_by(Course.created_at.desc()).limit(6).all()
    return template(request, 'admin/dashboard.html', db, {'admin': admin, 'stats': stats, 'recent_courses': recent_courses})


@app.get('/admin/login', response_class=HTMLResponse)
def admin_login_page(request: Request):
    with next_db_session() as db:
        return template(request, 'admin/login.html', db, {})


@app.post('/admin/login')
def admin_login(request: Request, username: str = Form(...), password: str = Form(...), db: Session = Depends(get_db)):
    admin = db.query(Admin).filter_by(username=username.strip()).first()
    if not admin or not verify_password(password, admin.password_hash):
        return RedirectResponse('/admin/login', status_code=303)
    request.session.clear()
    request.session['admin_id'] = admin.id
    return RedirectResponse('/admin/dashboard', status_code=303)


@app.get('/admin/logout')
def admin_logout(request: Request):
    request.session.clear()
    return RedirectResponse('/admin/login', status_code=303)


@app.get('/admin/programs')
def admin_programs(request: Request, db: Session = Depends(get_db)):
    admin = require_admin(request, db)
    programs = db.query(Program).order_by(Program.name).all()
    return template(request, 'admin/programs.html', db, {'admin': admin, 'programs': programs})


@app.post('/admin/programs')
def admin_create_program(request: Request, name: str = Form(...), description: str = Form(''), db: Session = Depends(get_db)):
    require_admin(request, db)
    db.add(Program(name=name.strip(), description=description.strip() or None))
    db.commit()
    return RedirectResponse('/admin/programs', status_code=303)


def cascade_delete_course(db: Session, course: Course):
    """Delete a course and every dependent row (lessons, materials, progress,
    quizzes, attempts, objectives, enrollments, purchases, certificates)."""
    lesson_ids = [row[0] for row in db.query(Lesson.id).filter(Lesson.course_id == course.id).all()]
    if lesson_ids:
        quiz_ids = [row[0] for row in db.query(Quiz.id).filter(Quiz.lesson_id.in_(lesson_ids)).all()]
        if quiz_ids:
            db.query(QuizAttempt).filter(QuizAttempt.quiz_id.in_(quiz_ids)).delete(synchronize_session=False)
            db.query(Quiz).filter(Quiz.id.in_(quiz_ids)).delete(synchronize_session=False)
        db.query(LessonMaterial).filter(LessonMaterial.lesson_id.in_(lesson_ids)).delete(synchronize_session=False)
        db.query(LessonProgress).filter(LessonProgress.lesson_id.in_(lesson_ids)).delete(synchronize_session=False)
        db.query(Lesson).filter(Lesson.id.in_(lesson_ids)).delete(synchronize_session=False)
    db.query(SessionObjective).filter(SessionObjective.course_id == course.id).delete(synchronize_session=False)
    db.query(Enrollment).filter(Enrollment.course_id == course.id).delete(synchronize_session=False)
    db.query(Purchase).filter(Purchase.course_id == course.id).delete(synchronize_session=False)
    db.query(CertificateAward).filter(CertificateAward.source_course_id == course.id).delete(synchronize_session=False)
    db.delete(course)


@app.post('/admin/programs/{program_id}/delete')
def admin_delete_program(program_id: int, request: Request, db: Session = Depends(get_db)):
    require_admin(request, db)
    program = db.get(Program, program_id)
    if program:
        try:
            for course in db.query(Course).filter(Course.program_id == program_id).all():
                cascade_delete_course(db, course)
            db.delete(program)
            db.commit()
        except Exception:
            db.rollback()
            logger.exception('Failed to delete program %s', program_id)
            return RedirectResponse('/admin/programs?error=1', status_code=303)
    return RedirectResponse('/admin/programs', status_code=303)


@app.get('/admin/courses')
def admin_courses(request: Request, db: Session = Depends(get_db)):
    admin = require_admin(request, db)
    courses = db.query(Course).order_by(Course.created_at.desc()).all()
    programs = db.query(Program).order_by(Program.name).all()
    return template(request, 'admin/courses.html', db, {'admin': admin, 'courses': courses, 'programs': programs})


@app.get('/admin/courses/new')
def admin_new_course(request: Request, db: Session = Depends(get_db)):
    admin = require_admin(request, db)
    programs = db.query(Program).order_by(Program.name).all()
    return template(request, 'admin/course_form.html', db, {'admin': admin, 'course': None, 'programs': programs, 'r2_ready': r2_enabled()})


@app.get('/admin/courses/{course_id}/edit')
def admin_edit_course(course_id: int, request: Request, db: Session = Depends(get_db)):
    admin = require_admin(request, db)
    course = db.get(Course, course_id)
    if not course:
        raise HTTPException(status_code=404)
    programs = db.query(Program).order_by(Program.name).all()
    return template(request, 'admin/course_form.html', db, {'admin': admin, 'course': course, 'programs': programs, 'r2_ready': r2_enabled()})


def price_cents(value):
    try:
        return max(0, int(round(float(value or 0) * 100)))
    except (TypeError, ValueError):
        return 0


@app.post('/admin/courses')
def admin_save_course(request: Request, course_id: int = Form(0), program_id: int = Form(...), title: str = Form(...),
                      description: str = Form(''), level: str = Form(''), slug: str = Form(''),
                      sales_copy: str = Form(''), thumbnail_url: str = Form(''), price: str = Form('0'),
                      currency: str = Form('USD'), num_lessons: int = Form(0), is_published: str = Form(''),
                      is_featured: str = Form(''), allow_free_enrollment: str = Form(''),
                      expertise_area: str = Form(''), certificate_level: int = Form(0), learning_hours: int = Form(0),
                      db: Session = Depends(get_db)):
    require_admin(request, db)
    course = db.get(Course, course_id) if course_id else Course(created_at=datetime.utcnow())
    course.program_id = program_id
    course.title = title.strip()
    course.description = description.strip() or None
    course.level = level.strip() or None
    course.slug = slug.strip() or slugify(course.title)
    course.sales_copy = sales_copy.strip() or None
    course.thumbnail_url = thumbnail_url.strip() or None
    course.price_cents = price_cents(price)
    course.currency = currency.strip().upper()[:3] or 'USD'
    course.expertise_area = expertise_area.strip() or None
    course.certificate_level = certificate_level if certificate_level in (0, 1, 2, 3) else 0
    if course.certificate_level and not learning_hours:
        learning_hours = CERTIFICATE_LEVEL_HOURS
    if course.certificate_level:
        num_lessons = MODULES_PER_LEVEL
    elif not num_lessons:
        num_lessons = MODULES_PER_LEVEL
    course.num_lessons = num_lessons
    course.learning_hours = max(0, learning_hours or 0)
    course.is_published = bool(is_published)
    course.is_featured = bool(is_featured)
    course.allow_free_enrollment = bool(allow_free_enrollment)
    is_new = not course_id
    if is_new:
        db.add(course)
        db.flush()
    blocked_extras = normalize_course_modules(db, course, num_lessons)
    if blocked_extras:
        db.rollback()
        return JSONResponse({
            'error': 'Extra lesson rows contain content and could not be removed automatically.',
            'extras': blocked_extras,
        }, status_code=409)
    db.commit()
    # Send a freshly-created course to its edit page so the admin can upload a
    # cover image and documents right away; edits go back to the course list.
    if is_new:
        return RedirectResponse(f'/admin/courses/{course.id}/edit?created=1', status_code=303)
    return RedirectResponse('/admin/courses', status_code=303)


@app.post('/admin/courses/{course_id}/delete')
def admin_delete_course(course_id: int, request: Request, db: Session = Depends(get_db)):
    require_admin(request, db)
    course = db.get(Course, course_id)
    if course:
        try:
            cascade_delete_course(db, course)
            db.commit()
        except Exception:
            db.rollback()
            logger.exception('Failed to delete course %s', course_id)
            return RedirectResponse('/admin/courses?error=1', status_code=303)
    return RedirectResponse('/admin/courses', status_code=303)


@app.post('/admin/lessons/{lesson_id}/clear')
def admin_clear_module(lesson_id: int, request: Request, db: Session = Depends(get_db)):
    """Wipe a module's content (materials, sessions, quizzes, progress) but keep the module slot."""
    require_admin(request, db)
    lesson = db.get(Lesson, lesson_id)
    if not lesson:
        raise HTTPException(status_code=404)
    module_number = current_module_number_for_lesson(db, lesson)
    quiz_ids = [row[0] for row in db.query(Quiz.id).filter(Quiz.lesson_id == lesson_id).all()]
    if quiz_ids:
        db.query(QuizAttempt).filter(QuizAttempt.quiz_id.in_(quiz_ids)).delete(synchronize_session=False)
        db.query(Quiz).filter(Quiz.id.in_(quiz_ids)).delete(synchronize_session=False)
    db.query(LessonMaterial).filter(LessonMaterial.lesson_id == lesson_id).delete(synchronize_session=False)
    db.query(LessonProgress).filter(LessonProgress.lesson_id == lesson_id).delete(synchronize_session=False)
    db.query(SessionObjective).filter(
        SessionObjective.course_id == lesson.course_id,
        SessionObjective.module_number == module_number,
    ).delete(synchronize_session=False)
    db.commit()
    return RedirectResponse(f'/admin/courses/{lesson.course_id}/lessons', status_code=303)


COURSE_DOC_FIELDS = {'syllabus': 'syllabus_file', 'clos': 'clos_file'}


@app.post('/admin/courses/{course_id}/documents/upload')
async def admin_upload_course_document(course_id: int, request: Request, doc_type: str = Form(...),
                                       file: UploadFile = File(...), db: Session = Depends(get_db)):
    """Upload a course-wide Syllabus or CLOs PDF to R2 (one per course)."""
    require_admin(request, db)
    course = db.get(Course, course_id)
    if not course:
        raise HTTPException(status_code=404)
    field = COURSE_DOC_FIELDS.get((doc_type or '').strip().lower())
    if not field:
        return RedirectResponse(f'/admin/courses/{course_id}/edit?doc=bad', status_code=303)
    if not r2_enabled():
        return RedirectResponse(f'/admin/courses/{course_id}/edit?doc=nor2', status_code=303)
    filename = (file.filename or '').strip()
    if not filename.lower().endswith('.pdf'):
        await file.close()
        return RedirectResponse(f'/admin/courses/{course_id}/edit?doc=pdf', status_code=303)
    content_type = file.content_type or guess_content_type(filename)
    key = object_key(filename, material_type=f'course-{field}')
    try:
        file.file.seek(0)
        upload_fileobj(key, file.file, content_type)
    except Exception as exc:
        logger.exception('Course document upload failed: %s', exc)
        return RedirectResponse(f'/admin/courses/{course_id}/edit?doc=fail', status_code=303)
    finally:
        await file.close()
    setattr(course, field, key)
    db.commit()
    return RedirectResponse(f'/admin/courses/{course_id}/edit?doc=ok', status_code=303)


@app.get('/courses/{course_id}/document/{kind}')
def course_document(course_id: int, kind: str, request: Request, db: Session = Depends(get_db)):
    """Serve a course-wide document (syllabus/clos) via a presigned R2 URL to admins or logged-in students."""
    field = COURSE_DOC_FIELDS.get((kind or '').strip().lower())
    if not field:
        raise HTTPException(status_code=404)
    course = db.get(Course, course_id)
    if not course:
        raise HTTPException(status_code=404)
    if not (admin_from_request(request, db) or student_from_request(request, db)):
        return RedirectResponse(f'/login?next=/courses/{course_id}/document/{kind}', status_code=303)
    key = getattr(course, field, None)
    if not key:
        raise HTTPException(status_code=404)
    label = 'Syllabus' if kind.strip().lower() == 'syllabus' else 'CLOs'
    return RedirectResponse(presigned_download_url(key, f'{label}-{course_slug(course)}.pdf'), status_code=303)


THUMB_EXTENSIONS = ('.jpg', '.jpeg', '.png', '.webp', '.gif')


@app.post('/admin/courses/{course_id}/thumbnail/upload')
async def admin_upload_course_thumbnail(course_id: int, request: Request,
                                        file: UploadFile = File(...), db: Session = Depends(get_db)):
    """Upload a course cover image (one per course) to R2; stored as an object key in thumbnail_url."""
    require_admin(request, db)
    course = db.get(Course, course_id)
    if not course:
        raise HTTPException(status_code=404)
    if not r2_enabled():
        return RedirectResponse(f'/admin/courses/{course_id}/edit?thumb=nor2', status_code=303)
    filename = (file.filename or '').strip()
    if not filename.lower().endswith(THUMB_EXTENSIONS):
        await file.close()
        return RedirectResponse(f'/admin/courses/{course_id}/edit?thumb=type', status_code=303)
    content_type = file.content_type or guess_content_type(filename)
    key = object_key(filename, material_type='course-thumbnail')
    try:
        file.file.seek(0)
        upload_fileobj(key, file.file, content_type)
    except Exception as exc:
        logger.exception('Course thumbnail upload failed: %s', exc)
        return RedirectResponse(f'/admin/courses/{course_id}/edit?thumb=fail', status_code=303)
    finally:
        await file.close()
    course.thumbnail_url = key
    db.commit()
    return RedirectResponse(f'/admin/courses/{course_id}/edit?thumb=ok', status_code=303)


@app.get('/courses/{course_id}/thumbnail')
def course_thumbnail(course_id: int, request: Request, db: Session = Depends(get_db)):
    """Serve a course cover image via a presigned R2 URL (public — cover art is not sensitive)."""
    course = db.get(Course, course_id)
    if not course or not course.thumbnail_url:
        raise HTTPException(status_code=404)
    key = course.thumbnail_url
    if key.startswith(('http://', 'https://', '/')):
        return RedirectResponse(key, status_code=303)
    return RedirectResponse(presigned_download_url(key, f'cover-{course_slug(course)}'), status_code=303)


@app.post('/admin/courses/{course_id}/bulk-prepare')
async def admin_prepare_course_bulk_import(course_id: int, request: Request, db: Session = Depends(get_db)):
    require_admin(request, db)
    course = db.get(Course, course_id)
    if not course:
        raise HTTPException(status_code=404)
    data = await request.json()
    module_names = [str(name).strip() for name in data.get('module_names', []) if str(name).strip()]
    if not module_names:
        return JSONResponse({'error': 'No module folders found.'}, status_code=400)
    module_names = module_names[:MODULES_PER_LEVEL]
    existing_lessons = db.query(Lesson).filter_by(course_id=course.id).order_by(Lesson.lesson_number).all()
    lesson_map = {}
    for index, module_name in enumerate(module_names, start=1):
        lesson = existing_lessons[index - 1] if index <= len(existing_lessons) else None
        if not lesson:
            lesson = Lesson(course_id=course.id, lesson_number=index, created_at=datetime.utcnow())
            db.add(lesson)
        lesson.lesson_number = index
        lesson.module_number = index
        lesson.session_number = 1
        lesson.duration_minutes = SESSION_DURATION_MINUTES
        lesson.title = re.sub(r'^\d+\s*[-_]\s*', '', module_name).replace('_', ' ').strip() or f'Module {index}'
        if not lesson.description:
            lesson.description = f'{lesson.title} module with five guided sessions, applied practice, and final simulation.'
        lesson_map[module_name] = lesson
    blocked_extras = []
    for extra in existing_lessons[len(module_names):]:
        has_content = (
            db.query(LessonMaterial).filter_by(lesson_id=extra.id).count()
            or db.query(LessonProgress).filter_by(lesson_id=extra.id).count()
            or db.query(Quiz).filter_by(lesson_id=extra.id).count()
        )
        if has_content:
            blocked_extras.append(extra.title)
        else:
            db.delete(extra)
    if blocked_extras:
        db.rollback()
        return JSONResponse({
            'error': 'This course has extra lesson rows with existing content. Remove them before bulk import.',
            'extras': blocked_extras,
        }, status_code=409)
    course.num_lessons = len(module_names)
    if not course.learning_hours:
        course.learning_hours = CERTIFICATE_LEVEL_HOURS
    db.commit()
    for lesson in lesson_map.values():
        db.refresh(lesson)
    return {
        'ok': True,
        'modules': [{'folder': folder, 'lesson_id': lesson.id, 'title': lesson.title} for folder, lesson in lesson_map.items()],
    }


@app.get('/admin/courses/{course_id}/lessons')
def admin_lessons(course_id: int, request: Request, db: Session = Depends(get_db)):
    admin = require_admin(request, db)
    course = db.get(Course, course_id)
    if not course:
        raise HTTPException(status_code=404)
    blocked_extras = normalize_course_modules(db, course, MODULES_PER_LEVEL)
    if not blocked_extras:
        db.commit()
    lessons = db.query(Lesson).filter_by(course_id=course_id).order_by(Lesson.lesson_number).all()
    return template(request, 'admin/lessons.html', db, {
        'admin': admin,
        'course': course,
        'lessons': lessons[:MODULES_PER_LEVEL],
        'blocked_extras': blocked_extras,
    })


@app.post('/admin/lessons/{lesson_id}')
def admin_update_lesson(lesson_id: int, request: Request, title: str = Form(...), description: str = Form(''), db: Session = Depends(get_db)):
    require_admin(request, db)
    lesson = db.get(Lesson, lesson_id)
    if not lesson:
        raise HTTPException(status_code=404)
    lesson.title = title.strip()
    lesson.description = description.strip() or None
    db.commit()
    return RedirectResponse(f'/admin/courses/{lesson.course_id}/lessons', status_code=303)


@app.get('/admin/lessons/{lesson_id}/materials')
def admin_materials(lesson_id: int, request: Request, db: Session = Depends(get_db)):
    admin = require_admin(request, db)
    lesson = db.get(Lesson, lesson_id)
    materials = db.query(LessonMaterial).filter_by(lesson_id=lesson_id).order_by(LessonMaterial.upload_order).all()
    current_module_number = current_module_number_for_lesson(db, lesson) if lesson else 1
    objective_rows = db.query(SessionObjective).filter_by(course_id=lesson.course_id, module_number=current_module_number).all() if lesson else []
    objective_by_session = {row.session_number: row for row in objective_rows}
    grouped = {key: [] for key, _label in MATERIAL_TYPE_OPTIONS}
    for material in materials:
        grouped.setdefault(material.material_type or 'other', []).append(material)
    return template(request, 'admin/materials.html', db, {
        'admin': admin,
        'lesson': lesson,
        'materials': materials,
        'grouped_materials': grouped,
        'material_type_options': MATERIAL_TYPE_OPTIONS,
        'r2_ready': r2_enabled(),
        'current_module_number': current_module_number,
        'objective_by_session': objective_by_session,
        'sessions_per_module': SESSIONS_PER_MODULE,
    })


@app.post('/admin/lessons/{lesson_id}/objectives')
async def admin_save_session_objectives(lesson_id: int, request: Request, db: Session = Depends(get_db)):
    require_admin(request, db)
    lesson = db.get(Lesson, lesson_id)
    if not lesson:
        raise HTTPException(status_code=404)
    current_module_number = current_module_number_for_lesson(db, lesson)
    form = await request.form()
    for session_number in range(1, SESSIONS_PER_MODULE + 1):
        title = (form.get(f'title_{session_number}') or '').strip()
        objective = (form.get(f'objective_{session_number}') or '').strip()
        existing = db.query(SessionObjective).filter_by(
            course_id=lesson.course_id,
            module_number=current_module_number,
            session_number=session_number,
        ).first()
        if not title and not objective:
            if existing:
                db.delete(existing)
            continue
        if not existing:
            existing = SessionObjective(
                course_id=lesson.course_id,
                module_number=current_module_number,
                session_number=session_number,
            )
            db.add(existing)
        existing.title = title or None
        existing.objective = objective or None
        existing.source = 'admin'
    db.commit()
    return RedirectResponse(f'/admin/lessons/{lesson.id}/materials', status_code=303)


@app.post('/admin/lessons/{lesson_id}/objectives/extract')
def admin_extract_session_objectives(lesson_id: int, request: Request, db: Session = Depends(get_db)):
    require_admin(request, db)
    lesson = db.get(Lesson, lesson_id)
    if not lesson:
        raise HTTPException(status_code=404)
    syllabus = db.query(LessonMaterial).filter_by(lesson_id=lesson.id, material_type='syllabus').order_by(LessonMaterial.upload_order.desc()).first()
    if not syllabus or not (syllabus.object_key or syllabus.file_path):
        return JSONResponse({'error': 'No syllabus PDF has been uploaded for this module.'}, status_code=400)
    try:
        text = extract_pdf_text(object_bytes(syllabus.object_key or syllabus.file_path))
        objectives = extract_topic_objectives(text)
    except Exception as exc:
        logger.exception('Syllabus objective extraction failed: %s', exc)
        return JSONResponse({'error': 'Could not extract objectives from the syllabus PDF.'}, status_code=500)
    if len(objectives) < SESSIONS_PER_MODULE:
        return JSONResponse({'error': f'Only {len(objectives)} objectives found in Topics covered.'}, status_code=422)
    module_number = current_module_number_for_lesson(db, lesson)
    saved = save_session_objectives(db, lesson.course_id, module_number, objectives)
    return {'ok': True, 'saved': saved, 'objectives': objectives[:SESSIONS_PER_MODULE]}


@app.post('/admin/lessons/{lesson_id}/materials/link')
async def admin_add_material_link(lesson_id: int, request: Request, material_type: str = Form('video'),
                                  video_url: str = Form(...), video_name: str = Form(''), db: Session = Depends(get_db)):
    require_admin(request, db)
    max_order = db.query(func.max(LessonMaterial.upload_order)).filter_by(lesson_id=lesson_id).scalar() or 0
    material = LessonMaterial(
        lesson_id=lesson_id,
        material_type=material_type,
        upload_order=max_order + 1,
        video_url=video_url.strip(),
        video_name=video_name.strip() or 'External resource',
        storage_provider='external',
    )
    db.add(material)
    db.commit()
    return RedirectResponse(f'/admin/lessons/{lesson_id}/materials', status_code=303)


@app.post('/admin/lessons/{lesson_id}/materials/presign')
async def admin_presign_material_upload(lesson_id: int, request: Request, db: Session = Depends(get_db)):
    require_admin(request, db)
    if not r2_enabled():
        return JSONResponse({'error': 'Cloudflare R2 is not configured.'}, status_code=400)
    data = await request.json()
    filename = (data.get('filename') or '').strip()
    content_type = (data.get('content_type') or guess_content_type(filename)).strip()
    if not filename:
        return JSONResponse({'error': 'filename is required'}, status_code=400)
    material_type = (data.get('material_type') or 'other').strip()
    key = object_key(filename, lesson_id=lesson_id, material_type=material_type)
    return {
        'upload_url': presigned_upload_url(key, content_type),
        'object_key': key,
        'content_type': content_type,
        'expires_in': cfg.R2_PRESIGN_EXPIRES_SECONDS,
    }


@app.post('/admin/lessons/{lesson_id}/materials/upload')
async def admin_upload_material_to_r2(lesson_id: int, request: Request, material_type: str = Form('other'),
                                      relative_path: str = Form(''), package_id: str = Form(''),
                                      register_metadata: str = Form('true'),
                                      file: UploadFile = File(...), db: Session = Depends(get_db)):
    require_admin(request, db)
    lesson = db.get(Lesson, lesson_id)
    if not lesson:
        raise HTTPException(status_code=404)
    if not r2_enabled():
        return JSONResponse({'error': 'Cloudflare R2 is not configured.'}, status_code=400)
    filename = (file.filename or '').strip()
    if not filename:
        return JSONResponse({'error': 'filename is required'}, status_code=400)

    clean_type = (material_type or 'other').strip() or 'other'
    content_type = file.content_type or guess_content_type(filename)
    clean_relative_path = (relative_path or filename).strip() or filename
    if package_id:
        key = package_object_key(clean_relative_path, lesson_id=lesson_id, material_type=clean_type, package_id=package_id)
    else:
        key = object_key(filename, lesson_id=lesson_id, material_type=clean_type)
    size_bytes = uploaded_file_size(file.file)

    try:
        file.file.seek(0)
        upload_fileobj(key, file.file, content_type)
    except Exception as exc:
        logger.exception('R2 material upload failed: %s', exc)
        return JSONResponse({'error': 'R2 upload failed before metadata was saved. Check R2 credentials, bucket permissions, and file size.'}, status_code=502)
    finally:
        await file.close()

    material = None
    if str(register_metadata).lower() not in ('false', '0', 'no', 'off'):
        try:
            material = save_r2_material_metadata(db, lesson_id, clean_type, filename, key, content_type, size_bytes)
        except Exception as exc:
            db.rollback()
            logger.exception('R2 metadata save failed after upload: %s', exc)
            return JSONResponse({
                'error': 'File reached R2, but the database record was not saved. Use Register uploaded file to repair it.',
                'object_key': key,
                'file_name': filename,
                'material_type': clean_type,
                'content_type': content_type,
                'size_bytes': size_bytes,
            }, status_code=500)
    return {
        'ok': True,
        'material_id': material.id if material else None,
        'file_name': material.file_name if material else filename,
        'material_type': material.material_type if material else clean_type,
        'object_key': material.object_key if material else key,
        'size_bytes': material.size_bytes if material else size_bytes,
    }


def save_r2_material_metadata(db: Session, lesson_id: int, material_type: str, file_name: str, object_key_value: str,
                              content_type: str = '', size_bytes: int = 0) -> LessonMaterial:
    existing = db.query(LessonMaterial).filter_by(lesson_id=lesson_id, object_key=object_key_value).first()
    if existing:
        existing.material_type = material_type
        existing.file_name = file_name or existing.file_name
        existing.file_path = object_key_value
        existing.content_type = content_type or existing.content_type
        existing.size_bytes = size_bytes or existing.size_bytes
        existing.storage_provider = 'r2'
        db.commit()
        db.refresh(existing)
        return existing
    max_order = db.query(func.max(LessonMaterial.upload_order)).filter_by(lesson_id=lesson_id).scalar() or 0
    material = LessonMaterial(
        lesson_id=lesson_id,
        material_type=material_type,
        file_name=file_name,
        file_path=object_key_value,
        object_key=object_key_value,
        content_type=content_type,
        size_bytes=size_bytes,
        storage_provider='r2',
        upload_order=max_order + 1,
    )
    db.add(material)
    db.commit()
    db.refresh(material)
    return material


@app.post('/admin/lessons/{lesson_id}/materials/register')
async def admin_register_uploaded_material(lesson_id: int, request: Request, db: Session = Depends(get_db)):
    require_admin(request, db)
    if not db.get(Lesson, lesson_id):
        raise HTTPException(status_code=404)
    data = await request.json()
    object_key_value = (data.get('object_key') or '').strip()
    if not object_key_value:
        return JSONResponse({'error': 'object_key is required'}, status_code=400)
    file_name = (data.get('file_name') or posixpath.basename(object_key_value)).strip()
    material_type = (data.get('material_type') or 'other').strip() or 'other'
    content_type = (data.get('content_type') or guess_content_type(file_name)).strip()
    size_bytes = int(data.get('size_bytes') or 0)
    material = save_r2_material_metadata(db, lesson_id, material_type, file_name, object_key_value, content_type, size_bytes)
    return {'ok': True, 'material_id': material.id}


@app.post('/admin/lessons/{lesson_id}/materials/sync-r2')
def admin_sync_lesson_r2_materials(lesson_id: int, request: Request, db: Session = Depends(get_db)):
    require_admin(request, db)
    if not db.get(Lesson, lesson_id):
        raise HTTPException(status_code=404)
    if not r2_enabled():
        return JSONResponse({'error': 'Cloudflare R2 is not configured.'}, status_code=400)
    existing_keys = {
        row[0] for row in db.query(LessonMaterial.object_key)
        .filter(LessonMaterial.lesson_id == lesson_id, LessonMaterial.object_key.isnot(None))
        .all()
    }
    created = 0
    prefix = f'lessons/{lesson_id}/'
    try:
        for item in list_objects(prefix):
            key = item['key']
            if not key or key.endswith('/') or key in existing_keys:
                continue
            file_name = posixpath.basename(key)
            material_type = infer_material_type_from_key(key)
            save_r2_material_metadata(
                db,
                lesson_id,
                material_type,
                file_name,
                key,
                guess_content_type(file_name),
                item.get('size') or 0,
            )
            existing_keys.add(key)
            created += 1
    except Exception as exc:
        logger.exception('R2 material sync failed: %s', exc)
        return JSONResponse({'error': 'Could not list lesson files in R2. Check R2 credentials and bucket permissions.'}, status_code=502)
    return RedirectResponse(f'/admin/lessons/{lesson_id}/materials?synced={created}', status_code=303)


@app.post('/admin/lessons/{lesson_id}/materials/complete')
async def admin_complete_material_upload(lesson_id: int, request: Request, db: Session = Depends(get_db)):
    require_admin(request, db)
    data = await request.json()
    max_order = db.query(func.max(LessonMaterial.upload_order)).filter_by(lesson_id=lesson_id).scalar() or 0
    material = LessonMaterial(
        lesson_id=lesson_id,
        material_type=(data.get('material_type') or 'article').strip(),
        file_name=(data.get('file_name') or '').strip(),
        file_path=(data.get('object_key') or '').strip(),
        object_key=(data.get('object_key') or '').strip(),
        content_type=(data.get('content_type') or '').strip(),
        size_bytes=int(data.get('size_bytes') or 0),
        storage_provider='r2',
        upload_order=max_order + 1,
    )
    if not material.object_key:
        return JSONResponse({'error': 'object_key is required'}, status_code=400)
    db.add(material)
    db.commit()
    return {'ok': True, 'material_id': material.id}


@app.post('/admin/materials/{material_id}/delete')
def admin_delete_material(material_id: int, request: Request, db: Session = Depends(get_db)):
    require_admin(request, db)
    material = db.get(LessonMaterial, material_id)
    lesson_id = material.lesson_id if material else 0
    if material:
        db.delete(material)
        db.commit()
    return RedirectResponse(f'/admin/lessons/{lesson_id}/materials', status_code=303)


@app.get('/materials/{material_id}/download')
def material_download(material_id: int, request: Request, db: Session = Depends(get_db)):
    material = db.get(LessonMaterial, material_id)
    if not material:
        raise HTTPException(status_code=404)
    if not material_access_allowed(material, request, db):
        return RedirectResponse(f'/login?next=/materials/{material_id}/download', status_code=303)
    if material.storage_provider == 'external' and material.video_url:
        return RedirectResponse(material.video_url, status_code=303)
    key = material.object_key or material.file_path
    if not key:
        raise HTTPException(status_code=404)
    if html_material(material):
        try:
            html = object_bytes(key).decode('utf-8')
        except UnicodeDecodeError:
            html = object_bytes(key).decode('utf-8', errors='replace')
        except Exception as exc:
            logger.exception('R2 HTML material load failed: %s', exc)
            raise HTTPException(status_code=502)
        return HTMLResponse(html)
    return RedirectResponse(presigned_download_url(key, material.file_name), status_code=303)


@app.get('/materials/{material_id}/{asset_path:path}')
def material_asset(material_id: int, asset_path: str, request: Request, db: Session = Depends(get_db)):
    material = db.get(LessonMaterial, material_id)
    if not material:
        raise HTTPException(status_code=404)
    if not material_access_allowed(material, request, db):
        return RedirectResponse(f'/login?next=/materials/{material_id}/download', status_code=303)
    if not html_material(material):
        raise HTTPException(status_code=404)
    base_key = posixpath.dirname(material.object_key or material.file_path or '')
    normalized_asset = posixpath.normpath(asset_path.replace('\\', '/')).lstrip('/')
    if normalized_asset.startswith('../') or normalized_asset == '..':
        raise HTTPException(status_code=404)
    key = posixpath.join(base_key, normalized_asset)
    filename = posixpath.basename(normalized_asset)
    return RedirectResponse(presigned_download_url(key, filename), status_code=303)


@app.get('/admin/companies')
def admin_companies(request: Request, db: Session = Depends(get_db)):
    admin = require_admin(request, db)
    companies = db.query(Company).order_by(Company.name).all()
    return template(request, 'admin/companies.html', db, {'admin': admin, 'companies': companies})


@app.post('/admin/companies')
def admin_save_company(request: Request, company_id: int = Form(0), name: str = Form(...), member_id: str = Form(''),
                       email: str = Form(''), is_active: str = Form('on'), db: Session = Depends(get_db)):
    require_admin(request, db)
    company = db.get(Company, company_id) if company_id else Company(created_at=datetime.utcnow())
    company.name = name.strip()
    company.member_id = member_id.strip() or None
    company.email = email.strip() or None
    company.is_active = bool(is_active)
    if not company_id:
        db.add(company)
    db.commit()
    return RedirectResponse('/admin/companies', status_code=303)


@app.get('/admin/students')
def admin_students(request: Request, db: Session = Depends(get_db)):
    admin = require_admin(request, db)
    students = db.query(Student).order_by(Student.created_at.desc()).all()
    companies = db.query(Company).order_by(Company.name).all()
    return template(request, 'admin/students.html', db, {'admin': admin, 'students': students, 'companies': companies})


@app.post('/admin/students')
def admin_save_student(request: Request, student_id: int = Form(0), full_name: str = Form(...), email: str = Form(...),
                       username: str = Form(''), password: str = Form(''), company_id: int = Form(0),
                       is_active: str = Form('on'), db: Session = Depends(get_db)):
    require_admin(request, db)
    student = db.get(Student, student_id) if student_id else Student(created_at=datetime.utcnow())
    student.full_name = full_name.strip()
    student.email = email.strip().lower()
    student.username = username.strip() or username_from_email(db, student.email)
    student.company_id = company_id or None
    student.is_active = bool(is_active)
    if password:
        student.password_hash = hash_password(password)
    elif not student_id:
        student.password_hash = hash_password('ChangeMe123!')
    if not student_id:
        db.add(student)
    db.commit()
    return RedirectResponse('/admin/students', status_code=303)


@app.get('/admin/enrollments')
def admin_enrollments(request: Request, db: Session = Depends(get_db)):
    admin = require_admin(request, db)
    enrollments = db.query(Enrollment).order_by(Enrollment.enrolled_at.desc()).all()
    students = db.query(Student).filter_by(is_active=True).order_by(Student.full_name).all()
    courses = db.query(Course).order_by(Course.title).all()
    return template(request, 'admin/enrollments.html', db, {'admin': admin, 'enrollments': enrollments, 'students': students, 'courses': courses})


@app.post('/admin/enrollments')
def admin_create_enrollment(request: Request, student_id: int = Form(...), course_id: int = Form(...), db: Session = Depends(get_db)):
    require_admin(request, db)
    enroll_student(db, student_id, course_id)
    db.commit()
    return RedirectResponse('/admin/enrollments', status_code=303)


@app.post('/admin/enrollments/{enrollment_id}/delete')
def admin_delete_enrollment(enrollment_id: int, request: Request, db: Session = Depends(get_db)):
    require_admin(request, db)
    enrollment = db.get(Enrollment, enrollment_id)
    if enrollment:
        db.delete(enrollment)
        db.commit()
    return RedirectResponse('/admin/enrollments', status_code=303)


@app.get('/admin/progression')
def admin_progression(request: Request, db: Session = Depends(get_db)):
    admin = require_admin(request, db)
    rows = []
    for enrollment in db.query(Enrollment).all():
        total = db.query(Lesson).filter_by(course_id=enrollment.course_id).count()
        done = db.query(LessonProgress).join(Lesson, LessonProgress.lesson_id == Lesson.id).filter(
            LessonProgress.student_id == enrollment.student_id,
            Lesson.course_id == enrollment.course_id,
            LessonProgress.is_completed.is_(True)
        ).count()
        rows.append({'enrollment': enrollment, 'done': done, 'total': total, 'pct': int(done / total * 100) if total else 0})
    return template(request, 'admin/progression.html', db, {'admin': admin, 'rows': rows})


@app.get('/admin/settings')
def admin_settings(request: Request, db: Session = Depends(get_db)):
    admin = require_admin(request, db)
    settings = {s.key: s.value for s in db.query(Settings).all()}
    return template(request, 'admin/settings.html', db, {'admin': admin, 'settings': settings})


@app.post('/admin/settings')
def admin_save_settings(request: Request, openai_api_key: str = Form(''), db: Session = Depends(get_db)):
    require_admin(request, db)
    setting = db.query(Settings).filter_by(key='openai_api_key').first()
    if not setting:
        setting = Settings(key='openai_api_key')
        db.add(setting)
    setting.value = openai_api_key.strip()
    db.commit()
    return RedirectResponse('/admin/settings', status_code=303)
